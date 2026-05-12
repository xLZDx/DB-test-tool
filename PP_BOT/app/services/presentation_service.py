"""Presentation service for PP_BOT.

Consumes architecture output and generates a PowerPoint deck with speaker notes
and slide visuals designed for stakeholder review.
"""
from __future__ import annotations

from pathlib import Path
import re
from typing import List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt

from ..config import PRESENTATION_DIR
from ..models import ArchitectureSection, PresentationRequest, PresentationResult, PresentationSlide
from ..prompts import PRESENTATION_SYSTEM_PROMPT


def _clean_lines(values: List[str]) -> List[str]:
    cleaned: List[str] = []
    seen: set[str] = set()
    for value in values:
        text = str(value).strip()
        if not text:
            continue
        key = text.lower()
        if key in seen:
            continue
        seen.add(key)
        cleaned.append(text)
    return cleaned


def _shorten(text: str, limit: int = 130) -> str:
    value = " ".join(str(text or "").split()).strip()
    if len(value) <= limit:
        return value
    return value[: limit - 3].rstrip() + "..."


def _slugify_filename(value: str) -> str:
    text = re.sub(r"[^A-Za-z0-9]+", "_", str(value or "").strip().lower()).strip("_")
    return text or "presentation"


class PresentationService:
    """Builds a stakeholder-friendly PowerPoint deck."""

    def _slide_from_section(
        self,
        section: ArchitectureSection,
        audience: str,
        include_notes: bool,
    ) -> PresentationSlide:
        subtitle = f"{audience.title()} view"
        notes = [section.body]
        if include_notes:
            notes.extend(section.bullets)
        image_prompts = [
            f"Create a clean visual that illustrates: {section.heading}",
            f"Use a professional enterprise style for {audience} stakeholders.",
        ]
        return PresentationSlide(
            title=section.heading,
            subtitle=subtitle,
            bullets=_clean_lines(section.bullets or [section.body]),
            notes=_clean_lines(notes),
            image_prompts=image_prompts,
        )

    def _build_outline(self, request: PresentationRequest) -> List[PresentationSlide]:
        architecture = request.architecture
        slides: List[PresentationSlide] = []

        slides.append(
            PresentationSlide(
                title=architecture.title,
                subtitle=request.topic,
                bullets=[
                    architecture.executive_summary,
                    f"Audience: {request.audience}",
                    f"Target slide count: {request.slide_count}",
                ],
                notes=[
                    architecture.executive_summary,
                    "Lead with the problem statement, then walk through the solution.",
                ],
                image_prompts=[
                    "Enterprise presentation cover visual",
                    "High-level workflow illustration for the topic",
                ],
            )
        )

        slides.append(
            PresentationSlide(
                title="Why this matters",
                subtitle="Business and technical drivers",
                bullets=_clean_lines(
                    architecture.sections[1].bullets if len(architecture.sections) > 1 else []
                ),
                notes=[
                    "Explain the business pain points and the technical constraints.",
                    "Frame the discussion around traceability, reliability, and reuse.",
                ],
                image_prompts=[
                    "A simple driver diagram showing business and technical constraints",
                ],
            )
        )

        for section in architecture.sections[2:]:
            slides.append(self._slide_from_section(section, request.audience, request.include_speaker_notes))

        slides.append(
            PresentationSlide(
                title="Key assumptions",
                subtitle="What we are assuming for the MVP",
                bullets=_clean_lines(architecture.assumptions or ["Initial MVP uses local fallback logic"]),
                notes=_clean_lines(
                    [
                        "Be explicit about assumptions so stakeholders know what is in and out of scope.",
                        *architecture.assumptions,
                    ]
                ),
                image_prompts=[
                    "A checklist or assumption callout visual",
                ],
            )
        )

        slides.append(
            PresentationSlide(
                title="Open questions",
                subtitle="Items to confirm before delivery",
                bullets=_clean_lines(
                    architecture.open_questions
                    or ["Confirm source systems", "Confirm branding and slide style"]
                ),
                notes=[
                    "Use this slide to capture unresolved integration and governance points.",
                ],
                image_prompts=[
                    "Question mark or decision tree visual",
                ],
            )
        )

        if len(slides) > request.slide_count:
            slides = slides[: request.slide_count]
        elif len(slides) < request.slide_count:
            while len(slides) < request.slide_count:
                slides.append(
                    PresentationSlide(
                        title="Implementation next step",
                        subtitle="Delivery planning",
                        bullets=[
                            "Finalize source connectors",
                            "Validate slide content with stakeholders",
                            "Refine visuals and narrative",
                        ],
                        notes=[
                            "Use filler slides only if the requested count is higher than the architecture outline.",
                        ],
                        image_prompts=[
                            "Roadmap or timeline visual",
                        ],
                    )
                )
        return slides

    def _set_wide_layout(self, prs: Presentation) -> None:
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    def _apply_branding(self, slide, brand_color: str, title: str, subtitle: Optional[str]) -> None:
        rgb = self._hex_to_rgb(brand_color)
        banner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.75))
        banner.fill.solid()
        banner.fill.fore_color.rgb = rgb
        banner.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.9), Inches(8.2), Inches(0.7))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(24)
        run.font.bold = True
        run.font.color.rgb = RGBColor(31, 31, 31)

        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(0.55), Inches(1.45), Inches(8.2), Inches(0.45))
            stf = subtitle_box.text_frame
            sp = stf.paragraphs[0]
            srun = sp.add_run()
            srun.text = subtitle
            srun.font.size = Pt(12)
            srun.font.color.rgb = RGBColor(95, 95, 95)

    def _add_bullets(self, slide, bullets: List[str]) -> None:
        box = slide.shapes.add_textbox(Inches(0.7), Inches(2.05), Inches(8.0), Inches(4.9))
        tf = box.text_frame
        tf.word_wrap = True
        first = True
        for bullet in bullets[:6]:
            paragraph = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            paragraph.text = _shorten(bullet, 180)
            paragraph.level = 0
            paragraph.font.size = Pt(18)

    def _add_visual_placeholder(self, slide, prompts: List[str], brand_color: str) -> None:
        rgb = self._hex_to_rgb(brand_color)
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(9.25),
            Inches(1.6),
            Inches(3.5),
            Inches(4.95),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = rgb
        box.fill.fore_color.brightness = 0.35
        box.line.color.rgb = rgb

        tf = box.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        r0 = p0.add_run()
        r0.text = "Visual / image cue"
        r0.font.size = Pt(16)
        r0.font.bold = True
        r0.font.color.rgb = RGBColor(255, 255, 255)

        for prompt in prompts[:3]:
            p = tf.add_paragraph()
            p.text = _shorten(prompt, 110)
            p.level = 0
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(255, 255, 255)

    def _add_notes(self, slide, notes: List[str]) -> None:
        try:
            notes_text = "\n".join(_clean_lines(notes))
            if notes_text:
                slide.notes_slide.notes_text_frame.text = notes_text
        except Exception:
            # Speaker notes are best-effort; deck generation should still succeed.
            pass

    def _hex_to_rgb(self, value: str) -> RGBColor:
        text = (value or "#1F4E79").strip().lstrip("#")
        if len(text) != 6:
            text = "1F4E79"
        return RGBColor(int(text[0:2], 16), int(text[2:4], 16), int(text[4:6], 16))

    def _write_deck(self, request: PresentationRequest, slides: List[PresentationSlide]) -> Path:
        PRESENTATION_DIR.mkdir(parents=True, exist_ok=True)
        output_path = PRESENTATION_DIR / f"{_slugify_filename(request.topic)}_deck.pptx"

        prs = Presentation()
        self._set_wide_layout(prs)

        for index, slide_data in enumerate(slides):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            self._apply_branding(slide, request.brand_color, slide_data.title, slide_data.subtitle)
            self._add_bullets(slide, slide_data.bullets)
            self._add_visual_placeholder(slide, slide_data.image_prompts, request.brand_color)
            if request.include_speaker_notes:
                self._add_notes(slide, slide_data.notes)

        prs.save(str(output_path))
        return output_path

    def generate(self, request: PresentationRequest) -> PresentationResult:
        slides = self._build_outline(request)
        output_path = self._write_deck(request, slides)
        return PresentationResult(
            topic=request.topic,
            title=request.architecture.title,
            subtitle=request.audience,
            slides=slides,
            output_path=str(output_path),
        )

    async def async_generate(self, request: PresentationRequest) -> PresentationResult:
        return self.generate(request)


presentation_service = PresentationService()
