from __future__ import annotations

import datetime as _dt
from dataclasses import dataclass
import pathlib
from typing import Iterable, Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_PATH = pathlib.Path("PP_BOT/data/outputs/CDS_CCAL_Transaction_Activity_Modern_Deck.pptx")

# Palette
NAVY = RGBColor(15, 23, 42)
SLATE = RGBColor(51, 65, 85)
MUTED = RGBColor(100, 116, 139)
TEXT = RGBColor(15, 23, 42)
WHITE = RGBColor(255, 255, 255)
SOFT_BG = RGBColor(248, 250, 252)
CARD = RGBColor(255, 255, 255)
BORDER = RGBColor(226, 232, 240)
BLUE = RGBColor(37, 99, 235)
CYAN = RGBColor(6, 182, 212)
GREEN = RGBColor(16, 185, 129)
AMBER = RGBColor(245, 158, 11)
ROSE = RGBColor(236, 72, 153)
INDIGO = RGBColor(99, 102, 241)
LIME = RGBColor(132, 204, 22)
ORANGE = RGBColor(249, 115, 22)


@dataclass
class SlideContext:
    prs: Presentation
    slide: any
    title: str
    subtitle: str
    page: int
    total: int
    dark: bool = False


def set_wide_layout(prs: Presentation) -> None:
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)


def rgb(value: str) -> RGBColor:
    value = value.strip().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def add_full_bg(slide, color: RGBColor) -> None:
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()


def add_footer(ctx: SlideContext) -> None:
    line = ctx.slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(7.14), Inches(13.333), Inches(0.36)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = rgb("EEF2FF")
    line.line.fill.background()

    left = ctx.slide.shapes.add_textbox(Inches(0.5), Inches(7.18), Inches(6.5), Inches(0.18))
    tf = left.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Source baseline: local brief, extracted diagrams, deck text, and SQL artifacts"
    r.font.name = "Aptos"
    r.font.size = Pt(9)
    r.font.color.rgb = MUTED

    right = ctx.slide.shapes.add_textbox(Inches(11.6), Inches(7.18), Inches(1.3), Inches(0.18))
    tf = right.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{ctx.page:02d}/{ctx.total:02d}"
    r.font.name = "Aptos"
    r.font.size = Pt(9)
    r.font.color.rgb = MUTED


def add_top_band(ctx: SlideContext, section: str) -> None:
    band = ctx.slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.48)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()

    accent = ctx.slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0.48), Inches(13.333), Inches(0.06)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = CYAN
    accent.line.fill.background()

    label = ctx.slide.shapes.add_textbox(Inches(0.45), Inches(0.08), Inches(8.8), Inches(0.22))
    tf = label.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = section
    r.font.name = "Aptos"
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = WHITE

    page = ctx.slide.shapes.add_textbox(Inches(11.9), Inches(0.08), Inches(1.0), Inches(0.22))
    tf = page.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{ctx.page:02d}/{ctx.total:02d}"
    r.font.name = "Aptos"
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = WHITE


def add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    font_size: int = 18,
    color: RGBColor = TEXT,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    font_name: str = "Aptos",
) -> any:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font_name
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_multiline_text(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    lines: Sequence[str],
    font_size: int = 16,
    color: RGBColor = TEXT,
    bullet: bool = False,
    font_name: str = "Aptos",
) -> any:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.clear()

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if bullet:
            p.level = 0
            p.bullet = True
        r = p.add_run()
        r.text = line
        r.font.name = font_name
        r.font.size = Pt(font_size)
        r.font.color.rgb = color
    return box


def add_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    body: Iterable[str] | str,
    accent: RGBColor = BLUE,
    fill: RGBColor = CARD,
    line: RGBColor = BORDER,
    title_size: int = 16,
    body_size: int = 12,
) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.3)

    stripe = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(0.12),
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()

    add_textbox(
        slide,
        left + 0.18,
        top + 0.18,
        width - 0.36,
        0.34,
        title,
        font_size=title_size,
        color=TEXT,
        bold=True,
    )

    if isinstance(body, str):
        body_lines = [body]
    else:
        body_lines = list(body)

    add_multiline_text(
        slide,
        left + 0.18,
        top + 0.60,
        width - 0.36,
        height - 0.75,
        body_lines,
        font_size=body_size,
        color=SLATE,
    )


def add_chip(
    slide,
    left: float,
    top: float,
    width: float,
    text: str,
    fill: RGBColor,
    color: RGBColor = WHITE,
    height: float = 0.34,
) -> None:
    chip = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    chip.fill.solid()
    chip.fill.fore_color.rgb = fill
    chip.line.fill.background()
    tf = chip.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = color


def add_arrow(slide, left: float, top: float, width: float, height: float, fill: RGBColor = BLUE) -> None:
    arrow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = fill
    arrow.line.fill.background()


def add_circle_node(
    slide,
    left: float,
    top: float,
    diameter: float,
    title: str,
    subtitle: str,
    fill: RGBColor,
    title_color: RGBColor = WHITE,
) -> None:
    node = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(left),
        Inches(top),
        Inches(diameter),
        Inches(diameter),
    )
    node.fill.solid()
    node.fill.fore_color.rgb = fill
    node.line.fill.background()
    tf = node.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = title
    r1.font.name = "Aptos"
    r1.font.size = Pt(14)
    r1.font.bold = True
    r1.font.color.rgb = title_color
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = subtitle
    r2.font.name = "Aptos"
    r2.font.size = Pt(9)
    r2.font.color.rgb = title_color


def add_note(slide, text: str) -> None:
    try:
        slide.notes_slide.notes_text_frame.text = text
    except Exception:
        pass


def create_title_slide(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    ctx = SlideContext(prs=prs, slide=slide, title="", subtitle="", page=page, total=total, dark=True)
    add_full_bg(slide, NAVY)

    # Accent geometry
    for left, top, width, height, color in [
        (9.25, 0.75, 3.0, 0.18, CYAN),
        (10.25, 1.05, 2.0, 0.18, GREEN),
        (8.85, 1.35, 3.4, 0.18, AMBER),
    ]:
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

    add_textbox(
        slide,
        0.62,
        0.42,
        5.8,
        0.35,
        "CDS / CCAL",
        font_size=14,
        color=CYAN,
        bold=True,
    )
    add_textbox(
        slide,
        0.62,
        1.05,
        6.8,
        1.20,
        "Transaction & Activity\nWorkflow",
        font_size=34,
        color=WHITE,
        bold=True,
    )
    add_textbox(
        slide,
        0.65,
        2.42,
        6.9,
        0.7,
        "A modern, diagram-rich deck covering transaction flow, activities, data model,\nschemas, tables, relationships, ODI jobs, and controls.",
        font_size=16,
        color=rgb("E2E8F0"),
    )

    # keyword chips
    chips = [
        "Transaction flow",
        "Transactions",
        "Activities",
        "CDS",
        "CCAL",
        "Tables",
        "Schemas",
        "Relationships",
        "E2E flow",
        "Transaction model",
    ]
    positions = [
        (0.65, 3.55, 1.55),
        (2.28, 3.55, 1.35),
        (3.72, 3.55, 1.25),
        (5.05, 3.55, 0.8),
        (5.92, 3.55, 0.9),
        (0.65, 3.98, 1.0),
        (1.75, 3.98, 1.0),
        (2.84, 3.98, 1.4),
        (4.36, 3.98, 1.15),
        (5.60, 3.98, 1.4),
    ]
    fills = [BLUE, CYAN, GREEN, AMBER, ROSE, INDIGO, RGBColor(14, 165, 233), RGBColor(20, 184, 166), RGBColor(168, 85, 247), ORANGE]
    for (label, (x, y, w), fill) in zip(chips, positions, fills):
        add_chip(slide, x, y, w, label, fill)

    # Right-side pipeline diagram card
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.45), Inches(0.82), Inches(4.0), Inches(5.75)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = rgb("0B1220")
    card.line.color.rgb = rgb("1E293B")

    add_textbox(slide, 8.8, 1.0, 2.8, 0.4, "E2E data path", font_size=16, color=CYAN, bold=True)

    steps = [
        ("Sources", "HPNS, trade servers,\nfiles, MQ, RJ Trust", CYAN),
        ("Landing", "Oracle landing tables", GREEN),
        ("Staging", "ODI jobs + validations", AMBER),
        ("Target", "CCAL_OWNER tables,\nextracts, consumers", ROSE),
    ]
    step_y = [1.55, 2.35, 3.15, 3.95]
    for (title, subtitle, fill), y in zip(steps, step_y):
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.9), Inches(y), Inches(2.9), Inches(0.58)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = fill
        box.line.fill.background()
        tf = box.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = f"{title}\n"
        r.font.name = "Aptos"
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = WHITE
        r2 = p.add_run()
        r2.text = subtitle
        r2.font.name = "Aptos"
        r2.font.size = Pt(9)
        r2.font.color.rgb = WHITE

    for y in [2.06, 2.86, 3.66]:
        add_arrow(slide, 10.0, y, 0.45, 0.22, fill=rgb("334155"))

    add_textbox(
        slide,
        8.86,
        4.9,
        3.0,
        0.55,
        "Purpose: show how transactions become activities,\npositions, balances, and audited outputs.",
        font_size=11,
        color=rgb("CBD5E1"),
    )

    add_textbox(
        slide,
        0.65,
        6.0,
        7.0,
        0.38,
        "Modern deck • diagrams built directly in PPTX • designed for executive + technical review",
        font_size=11,
        color=rgb("CBD5E1"),
    )

    add_note(
        slide,
        "Cover slide. Emphasize that CDS is the processing/control layer and CCAL is the business model for transactions, activities, positions, and balances."
    )
    add_footer(ctx)


def create_summary_slide(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    ctx = SlideContext(prs=prs, slide=slide, title="Executive summary", subtitle="", page=page, total=total)
    add_full_bg(slide, SOFT_BG)
    add_top_band(ctx, "Executive summary")

    add_textbox(slide, 0.55, 0.72, 7.0, 0.45, "What the CDS / CCAL platform does", font_size=24, bold=True)
    add_textbox(
        slide,
        0.55,
        1.15,
        8.9,
        0.45,
        "A transaction becomes an activity, the activity updates positions, balances and tax lots, and the platform keeps the full audit trail.",
        font_size=15,
        color=MUTED,
    )

    card_defs = [
        ("1  Source intake", "Multiple upstream feeds enter Oracle landing tables.\nExamples: retail, shadow, RJ Trust, RJ Bank, RJIG, files.", CYAN),
        ("2  Classification", "Transactions are normalized and classified into activities,\nsubtypes, and cash flow types.", GREEN),
        ("3  Persistence", "CDS and CCAL target schemas persist transaction, position,\nbalance, tax lot, and reference data.", AMBER),
        ("4  Control plane", "Exception pools, reference validation, and reconciliations keep\nthe flows traceable and operationally safe.", ROSE),
    ]
    xs = [0.55, 3.8, 7.05, 10.3]
    for (title, body, accent), x in zip(card_defs, xs):
        add_card(slide, x, 1.85, 2.7, 2.05, title, body.split("\n"), accent=accent, body_size=12)

    add_card(
        slide,
        0.55,
        4.25,
        12.2,
        1.55,
        "Key takeaways",
        [
            "CDS is the source-to-target orchestration layer.",
            "CCAL is the business layer for transactions, activities, positions, balances, and cost basis.",
            "ODI and GoldenGate carry the data through landing, staging, and target schemas.",
            "Controls are built around references, exceptions, duplicate checks, and reconciliation.",
        ],
        accent=BLUE,
        body_size=13,
    )

    add_note(
        slide,
        "Summarize the platform: multiple source families flow through CDS into CCAL, where transactions are classified, persisted, and rolled into position and balance outputs."
    )
    add_footer(ctx)


def create_e2e_flow_slide(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    ctx = SlideContext(prs=prs, slide=slide, title="E2E flow", subtitle="", page=page, total=total)
    add_full_bg(slide, WHITE)
    add_top_band(ctx, "End-to-end flow")

    add_textbox(slide, 0.55, 0.72, 8.5, 0.4, "Source -> Landing -> Staging -> Target -> Consumers", font_size=24, bold=True)

    # Main flow boxes
    labels = [
        ("Source systems", "HPNS\nTrade server\nRJ Trust\nRJ Bank\nShadow\nFiles", CYAN),
        ("Landing (Oracle)", "Raw intake tables\nSource-specific feeds", BLUE),
        ("CDS staging jobs", "ODI transforms\nReference lookup\nValidation\nException routing", GREEN),
        ("CCAL target tables", "TXN\nAPA\nFIP\nPOS\nBalances", AMBER),
        ("Consumers", "Reporting\nEOD extracts\nApplications\nDownstream services", ROSE),
    ]
    widths = [2.05, 2.05, 2.25, 2.05, 2.0]
    x = 0.55
    y = 1.75
    for (title, body, fill), w in zip(labels, widths):
        add_card(slide, x, y, w, 1.85, title, body.split("\n"), accent=fill, body_size=11)
        x += w + 0.25

    # Arrows between boxes
    x_positions = [2.60, 4.90, 7.40, 9.70]
    for ax in x_positions:
        add_arrow(slide, ax, 2.40, 0.22, 0.18, fill=rgb("CBD5E1"))

    # Side rails: reference and exceptions
    add_card(
        slide,
        0.55,
        4.15,
        5.95,
        1.6,
        "Reference data and validation path",
        [
            "REFERENCE_OWNER provides account and product reference enrichment.",
            "Reference lookups happen before target persistence and help classify flows.",
        ],
        accent=INDIGO,
        body_size=12,
    )
    add_card(
        slide,
        6.75,
        4.15,
        5.95,
        1.6,
        "Exception and audit path",
        [
            "CDS_UTIL_OWNER.EXCP_POOL and EXCP_DTL capture business and missing-reference issues.",
            "Ops analysts review failures and reconciliation exceptions.",
        ],
        accent=ROSE,
        body_size=12,
    )

    add_textbox(
        slide,
        0.58,
        6.0,
        12.0,
        0.34,
        "This is the repeatable pattern used across transaction, activity, position, and balance flows.",
        font_size=12,
        color=MUTED,
    )

    add_note(
        slide,
        "Explain the left-to-right pipeline and call out the two side paths: reference-data enrichment and exception handling. Mention landing, staging, target, and consumers."
    )
    add_footer(ctx)


def create_transaction_model_slide(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    ctx = SlideContext(prs=prs, slide=slide, title="Transaction model", subtitle="", page=page, total=total)
    add_full_bg(slide, SOFT_BG)
    add_top_band(ctx, "Transaction model and activity derivation")

    add_textbox(slide, 0.55, 0.72, 9.0, 0.42, "How a transaction becomes an activity and a position update", font_size=23, bold=True)

    # Core pipeline
    pipeline = [
        ("Transaction", "raw event from source", CYAN),
        ("Normalize", "standardize attributes\n& IDs", BLUE),
        ("Classify", "activity type\nsubtype\ncash flow", GREEN),
        ("Apply", "position update\nbalance effect\ncost basis", AMBER),
        ("Persist", "TXN / APA / FIP / POS", ROSE),
    ]
    px = [0.65, 2.95, 5.25, 7.55, 9.85]
    for (title, subtitle, fill), x in zip(pipeline, px):
        add_circle_node(slide, x, 1.95, 1.65, title, subtitle, fill)

    for ax in [2.48, 4.78, 7.08, 9.38]:
        add_arrow(slide, ax, 2.58, 0.26, 0.18, fill=rgb("CBD5E1"))

    add_card(
        slide,
        0.65,
        4.2,
        5.95,
        1.6,
        "Activity examples",
        [
            "Trades, transfers, ACAT, fees, dividends, corporate actions, checks, sweeps, withdrawals, deposits.",
            "The CCAL material references nearly 150 transaction types and ~1500 subtypes.",
        ],
        accent=GREEN,
        body_size=12,
    )
    add_card(
        slide,
        6.75,
        4.2,
        5.95,
        1.6,
        "What activity means",
        [
            "Activity is the client-facing and business-facing categorization of a transaction.",
            "It drives downstream reporting, cash-flow derivation, and reconciliation logic.",
        ],
        accent=AMBER,
        body_size=12,
    )

    add_textbox(
        slide,
        0.65,
        6.0,
        12.0,
        0.34,
        "Transaction -> activity -> position/balance update is the core CCAL processing model.",
        font_size=12,
        color=MUTED,
    )

    add_note(
        slide,
        "Describe the transaction model as a classification pipeline: normalize the raw event, classify it into an activity, and persist the effect into CCAL tables."
    )
    add_footer(ctx)


def create_schema_map_slide(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    ctx = SlideContext(prs=prs, slide=slide, title="Schemas and tables", subtitle="", page=page, total=total)
    add_full_bg(slide, WHITE)
    add_top_band(ctx, "Schemas and core tables")

    add_textbox(slide, 0.55, 0.72, 9.0, 0.42, "Core Oracle schemas and what each one owns", font_size=23, bold=True)

    # Center hub
    hub = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(4.85), Inches(2.15), Inches(3.2), Inches(2.0))
    hub.fill.solid()
    hub.fill.fore_color.rgb = NAVY
    hub.line.fill.background()
    tf = hub.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "CDS / CCAL\nProcessing Core"
    r.font.name = "Aptos"
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = WHITE

    hub_labels = [
        ("CCAL_OWNER", "TXN / APA / FIP / POS\nEOD_POS / AST_LBY_POS", 5.0, 0.95, BLUE),
        ("CCAL_REPL_OWNER", "SBDI_TX / SBDI_APA\nSBDI_FIP / ACTCCALQ_GG", 8.3, 1.75, CYAN),
        ("CDS_STG_OWNER", "TX_STG / CCAL_SUB_BAL_TGT\nAPPLICATION_PROPERTY", 8.3, 3.2, GREEN),
        ("CDS_UTIL_OWNER", "EXCP_POOL / EXCP_DTL", 5.0, 4.35, ROSE),
        ("REFERENCE_OWNER", "Account / Product reference", 1.35, 3.2, AMBER),
        ("Fill / source owners", "RJ Trust / RJ Bank / Shadow\nin-memory and feed-specific objects", 1.2, 1.0, INDIGO),
    ]
    for title, body, x, y, fill in hub_labels:
        add_card(slide, x, y, 2.95, 1.05, title, body.split("\n"), accent=fill, body_size=10, title_size=13)

    # Connector arrows toward the hub
    arrows = [
        (4.1, 1.35, 0.55, 0.18),
        (7.95, 1.9, 0.28, 0.55),
        (7.95, 3.3, 0.28, 0.55),
        (4.1, 4.6, 0.55, 0.18),
        (4.25, 3.0, 0.22, 0.18),
        (4.05, 2.2, 0.3, 0.18),
    ]
    for left, top, width, height in arrows:
        add_arrow(slide, left, top, width, height, fill=rgb("CBD5E1"))

    add_card(
        slide,
        0.55,
        5.55,
        12.2,
        0.95,
        "Important tables",
        [
            "POS, TXN, APA, FIP, TXN_STG, APA_STG, FIP_STG, CCAL_SUB_BAL_TGT, EXCP_POOL, EXCP_DTL, application_property",
        ],
        accent=INDIGO,
        body_size=11,
    )

    add_note(
        slide,
        "Walk through the schema map: CCAL_OWNER stores business entities, CCAL_REPL_OWNER handles replicated feeds, CDS_STG_OWNER carries staging and balance targets, CDS_UTIL_OWNER carries exceptions, and REFERENCE_OWNER supports lookup data."
    )
    add_footer(ctx)


def create_odi_slide(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    ctx = SlideContext(prs=prs, slide=slide, title="ODI orchestration", subtitle="", page=page, total=total)
    add_full_bg(slide, SOFT_BG)
    add_top_band(ctx, "ODI job orchestration and GoldenGate")

    add_textbox(slide, 0.55, 0.72, 9.0, 0.42, "How the jobs move data through the platform", font_size=23, bold=True)

    lanes = [
        ("Real-time lane", "REOCDS2EOD • FRNCDSBEOD • DIVCDS1BEOD • DIVCDS2BEOD • STRCCDSEOD • ACATCDSEOD", CYAN),
        ("Batch lane", "TRANCOMPLETE • S13CCALEOD • DIVACDSEOD • STRSCDSEOD", GREEN),
        ("Replication lane", "GoldenGate streams • ACTCCALQ_GG • STSCCALQ_GG", AMBER),
    ]
    top_positions = [1.55, 3.0, 4.45]
    for (title, body, fill), y in zip(lanes, top_positions):
        add_card(slide, 0.6, y, 12.1, 1.0, title, [body], accent=fill, body_size=12, title_size=14)

    # add a small timeline of processing stages
    stage_x = [0.95, 3.05, 5.15, 7.25, 9.35, 11.45]
    stage_titles = [
        ("Source", "feeds"),
        ("Landing", "Oracle"),
        ("Staging", "ODI"),
        ("Target", "CCAL_OWNER"),
        ("Extracts", "services"),
        ("Consumers", "reports"),
    ]
    for (title, subtitle), x in zip(stage_titles, stage_x):
        circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(6.0), Inches(0.55), Inches(0.55))
        circle.fill.solid()
        circle.fill.fore_color.rgb = NAVY if title in {"Landing", "Target"} else BLUE
        circle.line.fill.background()
        tf = circle.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = "•"
        r.font.name = "Aptos"
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = WHITE
        add_textbox(slide, x - 0.3, 6.58, 1.1, 0.28, f"{title}\n{subtitle}", font_size=10, color=SLATE, align=PP_ALIGN.CENTER)

    for ax in [1.63, 3.73, 5.83, 7.93, 10.03]:
        add_arrow(slide, ax, 6.18, 0.32, 0.16, fill=rgb("94A3B8"))

    add_note(
        slide,
        "Explain that ODI markers run either real-time or batch. Mention GoldenGate streams and the job names visible in the Visio extraction."
    )
    add_footer(ctx)


def create_position_balance_slide(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    ctx = SlideContext(prs=prs, slide=slide, title="Positions and balances", subtitle="", page=page, total=total)
    add_full_bg(slide, WHITE)
    add_top_band(ctx, "Position and balance processing")

    add_textbox(slide, 0.55, 0.72, 9.0, 0.42, "The position layer rolls into balances and end-of-day target data", font_size=23, bold=True)

    # flow diagram
    box_specs = [
        (0.7, 1.75, 2.0, 1.45, "Position events", "POS / EOD_POS\nAST_LBY_POS", CYAN),
        (3.1, 1.75, 2.05, 1.45, "Duplicate and date checks", "POS_DTM\nAST_LBY_POS_ID\nrow-count validation", AMBER),
        (5.55, 1.75, 2.1, 1.45, "Balance rollup", "CCAL_SUB_BAL_TGT\nCCAL_BAL_PBD_DT", GREEN),
        (8.05, 1.75, 2.0, 1.45, "Target and reconcile", "Oracle target table\nexception review", ROSE),
    ]
    for x, y, w, h, title, body, fill in box_specs:
        add_card(slide, x, y, w, h, title, body.split("\n"), accent=fill, body_size=11, title_size=13)

    for ax in [2.72, 5.15, 7.65]:
        add_arrow(slide, ax, 2.28, 0.28, 0.18, fill=rgb("CBD5E1"))

    add_card(
        slide,
        0.7,
        3.65,
        5.9,
        1.55,
        "What the balance SQL is doing",
        [
            "It inspects CCAL_OWNER.POS for specific AST_LBY_POS_ID values and date keys.",
            "It checks duplicates by source stream and POS_DTM before rollup.",
            "It updates CDS_STG_OWNER.APPLICATION_PROPERTY for CCAL_BAL_PBD_DT.",
        ],
        accent=BLUE,
        body_size=12,
    )
    add_card(
        slide,
        6.85,
        3.65,
        5.7,
        1.55,
        "Why that matters",
        [
            "Balances must be repeatable, auditable, and date-sensitive.",
            "Duplicate and exception handling protect downstream reporting integrity.",
        ],
        accent=INDIGO,
        body_size=12,
    )

    add_card(
        slide,
        0.7,
        5.55,
        12.0,
        0.95,
        "Tables and control points",
        [
            "POS, EOD_POS, AST_LBY_POS, CCAL_SUB_BAL_TGT, CCAL_BAL_STG_PART_EXC, application_property",
        ],
        accent=GREEN,
        body_size=11,
    )

    add_note(
        slide,
        "Explain the position and balance rollup. Mention the duplicate checks, the control property CCAL_BAL_PBD_DT, and how target rollups move into CCAL_SUB_BAL_TGT."
    )
    add_footer(ctx)


def create_controls_slide(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    ctx = SlideContext(prs=prs, slide=slide, title="Controls", subtitle="", page=page, total=total)
    add_full_bg(slide, SOFT_BG)
    add_top_band(ctx, "Controls, exceptions, and auditability")

    add_textbox(slide, 0.55, 0.72, 9.0, 0.42, "How the platform stays reliable and traceable", font_size=23, bold=True)

    control_cards = [
        ("Reference validation", "Account and product lookups occur before critical target writes.\nMissing refs route to exception handling.", CYAN),
        ("Exception pools", "CDS_UTIL_OWNER.EXCP_POOL and EXCP_DTL record failures,\nmissing refs, and operational issues.", ROSE),
        ("Reconciliation", "Transactions, positions, and balances are cross-checked\nfor duplicates, timing, and completeness.", GREEN),
    ]
    xs = [0.6, 4.55, 8.5]
    for (title, body, fill), x in zip(control_cards, xs):
        add_card(slide, x, 1.8, 3.35, 1.85, title, body.split("\n"), accent=fill, body_size=12, title_size=14)

    add_card(
        slide,
        0.6,
        4.0,
        12.15,
        1.35,
        "Operational review path",
        [
            "Ops analyst review of exceptions",
            "Reference data owner review",
            "Source system / feed rerun if needed",
            "Target refresh and reconciliation",
        ],
        accent=AMBER,
        body_size=12,
    )

    add_card(
        slide,
        0.6,
        5.55,
        12.15,
        0.95,
        "Audit principles",
        [
            "Trace every output back to a source feed, a job marker, and a controlled schema/table path.",
        ],
        accent=BLUE,
        body_size=11,
    )

    add_note(
        slide,
        "Highlight the exception path, the operational review loop, and how reference validation and reconciliation keep the platform stable."
    )
    add_footer(ctx)


def create_keywords_slide(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    ctx = SlideContext(prs=prs, slide=slide, title="Search terms", subtitle="", page=page, total=total)
    add_full_bg(slide, WHITE)
    add_top_band(ctx, "Search keywords and glossary")

    add_textbox(slide, 0.55, 0.72, 10.0, 0.42, "Search terms used to anchor the story and query the source material", font_size=23, bold=True)

    keywords = [
        ("Transaction flow", BLUE),
        ("Transactions", CYAN),
        ("Activities", GREEN),
        ("Activities flow", AMBER),
        ("CDS", ROSE),
        ("CCAL", INDIGO),
        ("Tables", ORANGE),
        ("Schemas", RGBColor(14, 165, 233)),
        ("Relationships", RGBColor(20, 184, 166)),
        ("E2E flow", RGBColor(168, 85, 247)),
        ("Transaction model", LIME),
        ("ODI", BLUE),
        ("GoldenGate", CYAN),
        ("Positions", GREEN),
        ("Balances", AMBER),
        ("Cost basis", ROSE),
        ("Tax lots", INDIGO),
        ("Oracle DB", ORANGE),
    ]
    x, y = 0.75, 1.85
    row_height = 0.48
    for word, fill in keywords:
        width = max(1.15, min(2.05, 0.42 + 0.12 * len(word)))
        if x + width > 12.55:
            x = 0.75
            y += 0.64
        add_chip(slide, x, y, width, word, fill)
        x += width + 0.18

    add_card(
        slide,
        0.75,
        4.8,
        12.0,
        1.4,
        "Mini glossary",
        [
            "CDS - control and integration layer that moves source data through landing, staging, and target.",
            "CCAL - business processing layer for transactions, activities, positions, balances, and cost basis.",
            "E2E - end-to-end path from source systems to reconciled business outputs.",
        ],
        accent=BLUE,
        body_size=12,
    )

    add_note(
        slide,
        "Use these search terms to orient future source searches or to explain the scope of the deck to stakeholders."
    )
    add_footer(ctx)


def create_closing_slide(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    ctx = SlideContext(prs=prs, slide=slide, title="Closing", subtitle="", page=page, total=total)
    add_full_bg(slide, NAVY)
    add_top_band(ctx, "Closing summary and next steps")

    add_textbox(slide, 0.65, 0.82, 9.0, 0.45, "What this deck gives you", font_size=23, bold=True, color=WHITE)
    add_textbox(
        slide,
        0.65,
        1.25,
        7.9,
        0.5,
        "A modern, diagram-heavy PPTX that explains CDS and CCAL transaction flow, activities, schemas, tables, controls, and balance rollup.",
        font_size=16,
        color=rgb("E2E8F0"),
    )

    roadmap = [
        ("1", "Review the diagram slides", CYAN),
        ("2", "Validate the component explanations", GREEN),
        ("3", "Replace placeholder visuals with official architecture images if available", AMBER),
        ("4", "Use the deck as the base for a longer design review", ROSE),
    ]
    y = 2.0
    for number, label, fill in roadmap:
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.8), Inches(y), Inches(0.5), Inches(0.5))
        dot.fill.solid()
        dot.fill.fore_color.rgb = fill
        dot.line.fill.background()
        tf = dot.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = number
        r.font.name = "Aptos"
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = WHITE

        add_textbox(slide, 1.45, y - 0.03, 8.5, 0.38, label, font_size=16, color=WHITE, bold=False)
        y += 0.62

    add_card(
        slide,
        8.8,
        1.78,
        3.8,
        2.4,
        "Output path",
        [
            "PP_BOT/data/outputs/CDS_CCAL_Transaction_Activity_Modern_Deck.pptx",
            "This deck is saved locally and can be opened in PowerPoint for review.",
        ],
        accent=CYAN,
        fill=rgb("0F172A"),
        line=rgb("1E293B"),
        body_size=11,
    )

    add_textbox(
        slide,
        0.65,
        6.0,
        11.2,
        0.34,
        "If you want, I can next turn this PPTX into a branded version with real source screenshots or add a PDF export step.",
        font_size=12,
        color=rgb("CBD5E1"),
    )

    add_note(
        slide,
        "Close by emphasizing the path from source systems to transactions, activities, positions, balances, and controls. Call out that the deck is editable and can be branded further."
    )
    add_footer(ctx)


def build_deck() -> pathlib.Path:
    prs = Presentation()
    set_wide_layout(prs)

    total = 10
    create_title_slide(prs, 1, total)
    create_summary_slide(prs, 2, total)
    create_e2e_flow_slide(prs, 3, total)
    create_transaction_model_slide(prs, 4, total)
    create_schema_map_slide(prs, 5, total)
    create_odi_slide(prs, 6, total)
    create_position_balance_slide(prs, 7, total)
    create_controls_slide(prs, 8, total)
    create_keywords_slide(prs, 9, total)
    create_closing_slide(prs, 10, total)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    return OUTPUT_PATH


def main() -> None:
    output = build_deck()
    print(f"Wrote {output}")


if __name__ == "__main__":
    main()
