from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_PATH = Path("PP_BOT/data/outputs/CDS_CCAL_Transaction_Activity_Only_Modern_Deck.pptx")

# Palette
NAVY = RGBColor(15, 23, 42)
SLATE = RGBColor(51, 65, 85)
MUTED = RGBColor(100, 116, 139)
TEXT = RGBColor(15, 23, 42)
WHITE = RGBColor(255, 255, 255)
SOFT_BG = RGBColor(247, 250, 252)
CARD = RGBColor(255, 255, 255)
BORDER = RGBColor(226, 232, 240)
BLUE = RGBColor(37, 99, 235)
CYAN = RGBColor(6, 182, 212)
GREEN = RGBColor(16, 185, 129)
AMBER = RGBColor(245, 158, 11)
ROSE = RGBColor(236, 72, 153)
INDIGO = RGBColor(99, 102, 241)
LIME = RGBColor(132, 204, 22)
ORANGE = RGBColor(249, 115, 22)
PALE_BLUE = RGBColor(219, 234, 254)
PALE_CYAN = RGBColor(207, 250, 254)
PALE_GREEN = RGBColor(220, 252, 231)
PALE_AMBER = RGBColor(254, 243, 199)


def rgb(value: str) -> RGBColor:
    value = value.strip().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def set_wide_layout(prs: Presentation) -> None:
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)


def add_full_bg(slide, color: RGBColor) -> None:
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()


def add_top_band(slide, section: str, page: int, total: int) -> None:
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.48))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()

    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0.48), Inches(13.333), Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = CYAN
    accent.line.fill.background()

    label = slide.shapes.add_textbox(Inches(0.45), Inches(0.08), Inches(8.8), Inches(0.22))
    tf = label.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = section
    r.font.name = "Aptos"
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = WHITE

    page_box = slide.shapes.add_textbox(Inches(11.9), Inches(0.08), Inches(1.0), Inches(0.22))
    tf = page_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{page:02d}/{total:02d}"
    r.font.name = "Aptos"
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = WHITE


def add_footer(slide) -> None:
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(7.14), Inches(13.333), Inches(0.36))
    line.fill.solid()
    line.fill.fore_color.rgb = rgb("EEF2FF")
    line.line.fill.background()

    left = slide.shapes.add_textbox(Inches(0.5), Inches(7.18), Inches(8.8), Inches(0.18))
    tf = left.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Source baseline: local Presentation folder, extracted diagrams, local SQL artifacts, and transaction/activity reference images"
    r.font.name = "Aptos"
    r.font.size = Pt(9)
    r.font.color.rgb = MUTED


def add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    font_size: int = 18,
    color: RGBColor = TEXT,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    font_name: str = "Aptos",
) -> any:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font_name
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_multiline(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    lines: Sequence[str],
    font_size: int = 14,
    color: RGBColor = TEXT,
    bullet: bool = False,
    font_name: str = "Aptos",
) -> any:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if bullet:
            p.bullet = True
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line
        r.font.name = font_name
        r.font.size = Pt(font_size)
        r.font.color.rgb = color
    return box


def add_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    body: Iterable[str] | str,
    accent: RGBColor = BLUE,
    fill: RGBColor = CARD,
    line: RGBColor = BORDER,
    title_size: int = 16,
    body_size: int = 12,
) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.2)

    stripe = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.10))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()

    add_textbox(slide, left + 0.18, top + 0.16, width - 0.36, 0.32, title, font_size=title_size, color=TEXT, bold=True)

    if isinstance(body, str):
        body_lines = [body]
    else:
        body_lines = list(body)

    add_multiline(slide, left + 0.18, top + 0.56, width - 0.36, height - 0.70, body_lines, font_size=body_size, color=SLATE)


def add_chip(slide, left: float, top: float, width: float, text: str, fill: RGBColor, color: RGBColor = WHITE) -> None:
    chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.34))
    chip.fill.solid()
    chip.fill.fore_color.rgb = fill
    chip.line.fill.background()
    tf = chip.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = color


def add_arrow(slide, left: float, top: float, width: float, height: float, fill: RGBColor = BLUE) -> None:
    arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(left), Inches(top), Inches(width), Inches(height))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = fill
    arrow.line.fill.background()


def add_line(slide, left: float, top: float, width: float, height: float, color: RGBColor = BORDER) -> None:
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def add_note(slide, text: str) -> None:
    try:
        slide.notes_slide.notes_text_frame.text = text
    except Exception:
        pass


def create_title_slide(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, NAVY)

    # decorative gradients/stripes
    for left, top, width, height, color in [
        (8.90, 0.75, 3.15, 0.18, CYAN),
        (9.95, 1.05, 2.10, 0.18, GREEN),
        (8.55, 1.35, 3.55, 0.18, AMBER),
    ]:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

    add_textbox(slide, 0.62, 0.42, 6.5, 0.36, "CDS / CCAL", font_size=14, color=CYAN, bold=True)
    add_textbox(slide, 0.62, 1.00, 7.0, 1.30, "Transactions &\nActivities Flow", font_size=34, color=WHITE, bold=True)
    add_textbox(
        slide,
        0.65,
        2.36,
        6.9,
        0.72,
        "Modern, diagram-first presentation for the transaction and activity lifecycle.\nFocused on source intake, classification, routing, tables, schemas, and relationships.",
        font_size=16,
        color=rgb("E2E8F0"),
    )

    chips = [
        ("Transaction flow", BLUE),
        ("Transactions", CYAN),
        ("Activities", GREEN),
        ("Activities flow", AMBER),
        ("CDS", ROSE),
        ("CCAL", INDIGO),
        ("Tables", ORANGE),
        ("Schemas", rgb("0EA5E9")),
        ("Relationships", rgb("14B8A6")),
        ("E2E flow", rgb("A855F7")),
        ("Transaction model", LIME),
    ]
    positions = [
        (0.65, 3.55, 1.55),
        (2.25, 3.55, 1.28),
        (3.60, 3.55, 1.10),
        (4.78, 3.55, 1.25),
        (6.12, 3.55, 0.80),
        (6.98, 3.55, 0.90),
        (0.65, 3.98, 1.00),
        (1.73, 3.98, 1.00),
        (2.81, 3.98, 1.30),
        (4.20, 3.98, 1.15),
        (5.43, 3.98, 1.40),
    ]
    for (label, fill), (x, y, w) in zip(chips, positions):
        add_chip(slide, x, y, w, label, fill)

    # right-side flow glyph
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.40), Inches(0.82), Inches(4.10), Inches(5.75))
    panel.fill.solid()
    panel.fill.fore_color.rgb = rgb("0B1220")
    panel.line.color.rgb = rgb("1E293B")

    add_textbox(slide, 8.78, 0.98, 2.80, 0.34, "Transaction activity path", font_size=16, color=CYAN, bold=True)

    for i, (title, subtitle, fill) in enumerate(
        [
            ("Sources", "Trade, file, MQ,\nQfile, direct feeds", CYAN),
            ("Landing", "Oracle intake", GREEN),
            ("Normalize", "Clean, map,\nstandardize", AMBER),
            ("Classify", "Activity and\ncash-flow rules", ROSE),
            ("Persist", "Tables & lineage", INDIGO),
        ]
    ):
        y = 1.45 + i * 0.72
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.85), Inches(y), Inches(2.95), Inches(0.52))
        box.fill.solid()
        box.fill.fore_color.rgb = fill
        box.line.fill.background()
        tf = box.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = f"{title}\n"
        r.font.name = "Aptos"
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = WHITE
        r2 = p.add_run()
        r2.text = subtitle
        r2.font.name = "Aptos"
        r2.font.size = Pt(9)
        r2.font.color.rgb = WHITE

    for y in [1.97, 2.69, 3.41, 4.13]:
        add_arrow(slide, 10.08, y, 0.46, 0.22, fill=rgb("334155"))

    add_textbox(
        slide,
        8.82,
        5.15,
        3.10,
        0.55,
        "Designed to read like a modern architecture presentation,\nnot a generic overview.",
        font_size=11,
        color=rgb("CBD5E1"),
    )

    add_note(
        slide,
        "Cover slide for the transaction/activity-only deck. Emphasize that this version excludes positions and balances and is focused on flow, processing, and tables."
    )
    add_footer(slide)


def create_scope_slide(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, SOFT_BG)
    add_top_band(slide, "Scope and narrative", page, total)

    add_textbox(slide, 0.55, 0.72, 8.0, 0.42, "What is included in this deck", font_size=23, bold=True)
    add_textbox(
        slide,
        0.55,
        1.15,
        10.8,
        0.48,
        "Transactions and activities only - source intake, e2e flow, classification, routing, tables, schemas, and exception handling.",
        font_size=15,
        color=MUTED,
    )

    cards = [
        ("1  Source families", "Retail, shadow, trust, bank, RJIG, files,\nQfile, direct feeds, corporate actions", CYAN),
        ("2  Transaction flow", "How the event moves from source -> landing -> staging -> target", BLUE),
        ("3  Activity model", "How a transaction is normalized and classified into an activity", GREEN),
        ("4  Controls", "Validation, exceptions, audit trail,\nreconciliation, operational review", ROSE),
    ]
    xs = [0.55, 3.75, 6.95, 10.15]
    for (title, body, accent), x in zip(cards, xs):
        add_card(slide, x, 1.85, 2.65, 2.05, title, body.split("\n"), accent=accent, body_size=12)

    add_card(
        slide,
        0.55,
        4.25,
        12.2,
        1.45,
        "How this presentation is structured",
        [
            "Start with the transaction source landscape.",
            "Show the end-to-end flow with a diagram translated from the local reference style.",
            "Explain the activity model, schemas, table relationships, and processing jobs.",
            "Close with control points and the key search terms used in the design.",
        ],
        accent=INDIGO,
        body_size=13,
    )

    add_note(
        slide,
        "This slide sets the scope: transactions and activities only. Tell the audience that positions and balances are intentionally out of scope."
    )
    add_footer(slide)


def create_source_landscape_slide(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, WHITE)
    add_top_band(slide, "Transaction source landscape", page, total)

    add_textbox(slide, 0.55, 0.72, 8.5, 0.42, "Transaction sources and how they feed the flow", font_size=23, bold=True)

    source_groups = [
        ("Real-time feeds", ["Retail/PCG", "Shadow", "Standard Trade Server"], CYAN, PALE_CYAN),
        ("Bank / trust", ["RJ Bank", "RJ Trust", "JDBC / in-memory"], BLUE, PALE_BLUE),
        ("Market / institutional", ["RJIG", "Impact", "Corporate actions"], GREEN, PALE_GREEN),
        ("Files / special", ["Qfile", "Files", "Direct mutual funds"], AMBER, PALE_AMBER),
    ]
    x_positions = [0.55, 3.75, 6.95, 10.15]
    for (title, items, accent, fill), x in zip(source_groups, x_positions):
        add_card(slide, x, 1.80, 2.65, 2.4, title, items, accent=accent, fill=fill, body_size=13, title_size=14)

    add_card(
        slide,
        0.55,
        4.55,
        12.2,
        1.25,
        "Common intake pattern",
        [
            "Upstream feeds land in Oracle staging structures, then route into transaction normalization and activity classification.",
            "Some sources are real-time, others are batch or end-of-day, and some need reclassification or manual review.",
        ],
        accent=INDIGO,
        body_size=12,
    )

    add_textbox(
        slide,
        0.55,
        6.00,
        12.0,
        0.30,
        "Use the source landscape to explain why different transaction families follow different control paths.",
        font_size=12,
        color=MUTED,
    )

    add_note(
        slide,
        "Introduce the different source families. Explain that the flow is not one pipeline; it is a set of source-specific branches that converge into common transaction processing."
    )
    add_footer(slide)


def create_reference_style_flow_slide(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, SOFT_BG)
    add_top_band(slide, "Translated flow diagram", page, total)

    add_textbox(slide, 0.55, 0.72, 10.5, 0.42, "Transaction flow translated from the local reference image", font_size=23, bold=True)
    add_textbox(
        slide,
        0.55,
        1.10,
        10.5,
        0.30,
        "A wide diagram that mirrors the local style: many source boxes at the top, common intake in the middle, and activity outcomes below.",
        font_size=13,
        color=MUTED,
    )

    # Top lanes of source boxes
    labels = [
        "Retail / PCG",
        "Shadow",
        "RJ Trust",
        "RJ Bank",
        "RJIG",
        "Impact",
        "Files",
        "Qfile",
        "ACATS",
        "Corporate actions",
        "Direct funds",
        "Special reclass",
    ]
    colors = [CYAN, BLUE, GREEN, AMBER, ROSE, INDIGO, ORANGE, rgb("0EA5E9"), rgb("14B8A6"), rgb("A855F7"), LIME, rgb("14B8A6")]
    xs = [0.58 + i * 1.02 for i in range(12)]

    # top row boxes
    for i, (label, fill, x) in enumerate(zip(labels, colors, xs)):
        add_chip(slide, x, 1.62, 0.92, label, fill)

    # small second row boxes
    second_row = [
        "Trade",
        "Transfer",
        "Trust txn",
        "Bank txn",
        "Market evt",
        "Impact",
        "File load",
        "Qfile load",
        "ACAT",
        "Corp action",
        "Direct",
        "Reclass",
    ]
    for i, (label, x) in enumerate(zip(second_row, xs)):
        add_card(slide, x - 0.03, 2.10, 0.98, 0.52, "", [label], accent=colors[i], body_size=8, title_size=8, fill=WHITE)

    # vertical connectors from top to second row and a common base line
    for x in [0.99 + i * 1.02 for i in range(12)]:
        add_line(slide, x, 1.96, 0.02, 0.12, color=rgb("C7D2FE"))
        add_line(slide, x, 2.62, 0.02, 1.48, color=rgb("CBD5E1"))

    add_line(slide, 0.95, 4.10, 11.4, 0.03, color=rgb("94A3B8"))

    central = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(4.55), Inches(2.76), Inches(4.20), Inches(0.94))
    central.fill.solid()
    central.fill.fore_color.rgb = NAVY
    central.line.fill.background()
    tf = central.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Transaction normalization and activity routing"
    r.font.name = "Aptos"
    r.font.size = Pt(17)
    r.font.bold = True
    r.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "common flow / common controls / traceability"
    r2.font.name = "Aptos"
    r2.font.size = Pt(10)
    r2.font.color.rgb = rgb("CBD5E1")

    # lower branches
    lower_cards = [
        ("Cash flow derivation", "contributions\nwithdrawals\nsecurities movement", CYAN, 1.45),
        ("Activity reclassification", "history / correction\nmanual review\nrule refresh", AMBER, 5.05),
        ("Exception review", "missing reference\nduplicates\nrerun queue", ROSE, 8.65),
    ]
    for title, body, accent, x in lower_cards:
        add_card(slide, x, 4.55, 3.0, 1.45, title, body.split("\n"), accent=accent, body_size=12)

    # lines to lower cards
    for x in [2.95, 6.45, 10.05]:
        add_arrow(slide, x, 3.82, 0.38, 0.20, fill=rgb("CBD5E1"))

    add_textbox(
        slide,
        0.55,
        6.20,
        12.0,
        0.24,
        "This slide intentionally mirrors the local diagram style: many small source boxes, a shared intake line, and a lower activity/output layer.",
        font_size=11,
        color=MUTED,
    )

    add_note(
        slide,
        "Use this as the main visual reference slide. Emphasize that the original image had a wide horizontal topology and the deck now translates that into clean PPT shapes."
    )
    add_footer(slide)


def create_activity_model_slide(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, WHITE)
    add_top_band(slide, "Activity model and transaction processing", page, total)

    add_textbox(slide, 0.55, 0.72, 9.0, 0.42, "How a transaction becomes an activity", font_size=23, bold=True)

    steps = [
        ("Transaction event", "raw source record", CYAN),
        ("Normalize", "ids, dates,\nvalues, source fields", BLUE),
        ("Classify", "activity type\nsubtype\ncash flow category", GREEN),
        ("Route", "business rules\nline of business\nexceptions", AMBER),
        ("Persist", "transaction tables\nactivity history\naudit trail", ROSE),
    ]
    xs = [0.65, 2.85, 5.05, 7.25, 9.45]
    for (title, subtitle, fill), x in zip(steps, xs):
        node = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.00), Inches(1.75), Inches(1.25))
        node.fill.solid()
        node.fill.fore_color.rgb = fill
        node.line.fill.background()
        tf = node.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = title + "\n"
        r.font.name = "Aptos"
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = WHITE
        r2 = p.add_run()
        r2.text = subtitle
        r2.font.name = "Aptos"
        r2.font.size = Pt(9)
        r2.font.color.rgb = WHITE

    for x in [2.52, 4.72, 6.92, 9.12]:
        add_arrow(slide, x, 2.45, 0.28, 0.18, fill=rgb("CBD5E1"))

    add_card(
        slide,
        0.75,
        4.10,
        3.85,
        1.55,
        "Activity examples",
        [
            "Trades, transfers, fees, ACAT, dividends, sweeps, checks, corporate actions, direct investments.",
            "The model supports dozens of activity families and many more subtype refinements.",
        ],
        accent=GREEN,
        body_size=12,
    )
    add_card(
        slide,
        4.80,
        4.10,
        3.85,
        1.55,
        "What activities do",
        [
            "They normalize business meaning out of raw transactions.",
            "They drive routing, downstream reporting, and exception handling.",
        ],
        accent=AMBER,
        body_size=12,
    )
    add_card(
        slide,
        8.85,
        4.10,
        3.85,
        1.55,
        "Why this matters",
        [
            "Activities are the stable business layer above raw transaction events.",
            "The same flow works across source families and operational windows.",
        ],
        accent=ROSE,
        body_size=12,
    )

    # activity taxonomy chips
    taxonomy = [
        "Trade",
        "Transfer",
        "ACAT",
        "Dividend",
        "Fee",
        "Corporate action",
        "Check",
        "Sweep",
        "Cash deposit",
        "Cash withdrawal",
        "Direct mutual fund",
        "Special reclass",
    ]
    x = 0.75
    y = 6.00
    for idx, label in enumerate(taxonomy):
        width = max(0.95, min(1.65, 0.35 + 0.11 * len(label)))
        if x + width > 12.55:
            x = 0.75
            y += 0.42
        fill = [CYAN, BLUE, GREEN, AMBER, ROSE, INDIGO][idx % 6]
        add_chip(slide, x, y, width, label, fill)
        x += width + 0.12

    add_note(
        slide,
        "Walk through the transaction-to-activity model. Each arrow represents normalization, classification, routing, and persistence."
    )
    add_footer(slide)


def create_schema_relationship_slide(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, SOFT_BG)
    add_top_band(slide, "Schemas, tables, and relationships", page, total)

    add_textbox(slide, 0.55, 0.72, 9.4, 0.42, "Transaction-centric schema map", font_size=23, bold=True)

    hub = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(4.70), Inches(2.15), Inches(3.35), Inches(2.05))
    hub.fill.solid()
    hub.fill.fore_color.rgb = NAVY
    hub.line.fill.background()
    tf = hub.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "CDS / CCAL\ntransaction core"
    r.font.name = "Aptos"
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = WHITE

    nodes = [
        ("CCAL_OWNER", "TXN / TXN_STG\nAPA / FIP", 5.00, 0.95, BLUE),
        ("CCAL_REPL_OWNER", "ACTCCALQ_GG\nSBDI_TX", 8.25, 1.75, CYAN),
        ("CDS_STG_OWNER", "TXN_STG\nAPPLICATION_PROPERTY", 8.25, 3.25, GREEN),
        ("CDS_UTIL_OWNER", "EXCP_POOL\nEXCP_DTL", 4.90, 4.40, ROSE),
        ("REFERENCE_OWNER", "account / product\nlookup tables", 1.20, 3.25, AMBER),
        ("Source owners", "RJ Trust / RJ Bank\nShadow / file feeds", 1.05, 1.00, INDIGO),
    ]
    for title, body, x, y, fill in nodes:
        add_card(slide, x, y, 3.00, 1.05, title, body.split("\n"), accent=fill, body_size=10, title_size=13)

    # arrows toward center hub
    for left, top, width, height in [
        (4.15, 1.35, 0.50, 0.18),
        (7.95, 1.90, 0.24, 0.55),
        (7.95, 3.32, 0.24, 0.55),
        (4.25, 4.65, 0.50, 0.18),
        (4.20, 2.25, 0.26, 0.18),
    ]:
        add_arrow(slide, left, top, width, height, fill=rgb("CBD5E1"))

    add_card(
        slide,
        0.70,
        5.55,
        12.0,
        0.98,
        "Primary tables",
        [
            "Transaction tables, staging tables, exception tables, reference tables, and activity history tables",
        ],
        accent=INDIGO,
        body_size=11,
    )

    add_note(
        slide,
        "Explain the schema relationships around transaction processing. Keep it centered on transaction tables, staging, reference data, and exception handling."
    )
    add_footer(slide)


def create_qfile_slide(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, WHITE)
    add_top_band(slide, "Qfile and special transaction branches", page, total)

    add_textbox(slide, 0.55, 0.72, 9.2, 0.42, "Special branches that feed the transaction activity flow", font_size=23, bold=True)

    groups = [
        ("Qfile sourcing", ["transaction qfile", "edited flow", "special source feed"], CYAN, PALE_CYAN),
        ("Manual / special", ["reclassification", "re-establish failed txns", "review queue"], AMBER, PALE_AMBER),
        ("Corporate actions", ["dividends", "splits", "mergers", "spin-offs"], GREEN, PALE_GREEN),
        ("Cash movement", ["deposits", "withdrawals", "checks", "sweeps"], ROSE, rgb("FCE7F3")),
    ]
    xs = [0.55, 3.75, 6.95, 10.15]
    for (title, items, accent, fill), x in zip(groups, xs):
        add_card(slide, x, 1.78, 2.65, 2.28, title, items, accent=accent, fill=fill, body_size=12, title_size=14)

    add_card(
        slide,
        0.55,
        4.40,
        12.2,
        1.65,
        "How the branch behaves",
        [
            "A Qfile or special feed may not look like a standard trade stream, but it still becomes a normalized transaction event.",
            "After classification it is routed to the same activity model and stored with the same lineage and audit expectations.",
        ],
        accent=BLUE,
        body_size=12,
    )

    # mini flow bar
    add_line(slide, 1.05, 6.28, 10.9, 0.03, color=rgb("CBD5E1"))
    mini = [
        ("Qfile / file", 1.00, CYAN),
        ("Normalize", 3.30, BLUE),
        ("Classify activity", 5.65, GREEN),
        ("Route / exception", 8.00, AMBER),
        ("Persist", 10.40, ROSE),
    ]
    for label, x, fill in mini:
        node = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(5.96), Inches(1.60), Inches(0.42))
        node.fill.solid()
        node.fill.fore_color.rgb = fill
        node.line.fill.background()
        tf = node.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.name = "Aptos"
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.color.rgb = WHITE

    add_note(
        slide,
        "Use this slide for the special branches: Qfile, file sourcing, and reclassification. It reinforces that all paths converge on the same transaction/activity model."
    )
    add_footer(slide)


def create_orchestration_slide(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, SOFT_BG)
    add_top_band(slide, "ODI orchestration and processing markers", page, total)

    add_textbox(slide, 0.55, 0.72, 9.2, 0.42, "Job markers and lanes used in transaction processing", font_size=23, bold=True)

    lanes = [
        ("Real-time lane", "REOCDS2EOD • FRNCDSBEOD • DIVCDS1BEOD • DIVCDS2BEOD • STRCCDSEOD • ACATCDSEOD", CYAN),
        ("Batch lane", "TRANCOMPLETE • S13CCALEOD • DIVACDSEOD • STRSCDSEOD", GREEN),
        ("Special handling lane", "Qfile / file rerun • manual reclass • reference validation", AMBER),
    ]
    ys = [1.55, 2.95, 4.35]
    for (title, body, fill), y in zip(lanes, ys):
        add_card(slide, 0.6, y, 12.1, 1.02, title, [body], accent=fill, body_size=12, title_size=14)

    # timeline
    labels = ["Source", "Landing", "Staging", "Classify", "Target", "Extracts"]
    fill_colors = [BLUE, CYAN, GREEN, AMBER, ROSE, INDIGO]
    x_positions = [0.95, 3.05, 5.15, 7.25, 9.35, 11.45]
    for (label, subtitle, fill), x in zip(zip(labels, ["feeds", "Oracle", "ODI", "rules", "tables", "consumers"], fill_colors), x_positions):
        circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(6.0), Inches(0.55), Inches(0.55))
        circle.fill.solid()
        circle.fill.fore_color.rgb = fill
        circle.line.fill.background()
        tf = circle.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = "•"
        r.font.name = "Aptos"
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = WHITE
        add_textbox(slide, x - 0.3, 6.58, 1.1, 0.28, f"{label}\n{subtitle}", font_size=10, color=SLATE, align=PP_ALIGN.CENTER)

    for ax in [1.63, 3.73, 5.83, 7.93, 10.03]:
        add_arrow(slide, ax, 6.18, 0.32, 0.16, fill=rgb("94A3B8"))

    add_note(
        slide,
        "Describe the processing lanes and markers. Keep the narrative centered on orchestration of transaction and activity flows, not positions or balances."
    )
    add_footer(slide)


def create_controls_slide(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, WHITE)
    add_top_band(slide, "Controls and auditability", page, total)

    add_textbox(slide, 0.55, 0.72, 9.0, 0.42, "How the transaction flow stays trusted", font_size=23, bold=True)

    cards = [
        ("Reference validation", "Validate account/product references before critical target writes", CYAN),
        ("Duplicate detection", "Use source + transaction keys to prevent duplicate events and reruns", AMBER),
        ("Exception pools", "Route missing-reference and business-rule failures to EXCP_POOL / EXCP_DTL", ROSE),
        ("Audit trail", "Preserve lineage from source feed to activity classification and persisted output", GREEN),
    ]
    xs = [0.55, 3.70, 6.85, 10.00]
    for (title, body, accent), x in zip(cards, xs):
        add_card(slide, x, 1.75, 2.75, 1.95, title, [body], accent=accent, body_size=11, title_size=14)

    add_card(
        slide,
        0.55,
        4.10,
        12.2,
        1.45,
        "Operational review path",
        [
            "Ops analyst review",
            "Reference owner review",
            "Source rerun or correction",
            "Target refresh and reconciliation",
        ],
        accent=INDIGO,
        body_size=12,
    )

    add_card(
        slide,
        0.55,
        5.75,
        12.2,
        0.75,
        "Lineage principle",
        [
            "Every activity should be traceable back to a source, a marker/job, and a controlled schema/table path.",
        ],
        accent=BLUE,
        body_size=11,
    )

    add_note(
        slide,
        "Close the process explanation by emphasizing reference validation, duplicate detection, exception handling, and the audit trail."
    )
    add_footer(slide)


def create_keywords_slide(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, SOFT_BG)
    add_top_band(slide, "Search terms and glossary", page, total)

    add_textbox(slide, 0.55, 0.72, 9.8, 0.42, "Keywords used to frame the deck and search the local material", font_size=23, bold=True)

    words = [
        ("Transaction flow", BLUE),
        ("Transactions", CYAN),
        ("Activities", GREEN),
        ("Activities flow", AMBER),
        ("CDS", ROSE),
        ("CCAL", INDIGO),
        ("Tables", ORANGE),
        ("Schemas", rgb("0EA5E9")),
        ("Relationships", rgb("14B8A6")),
        ("E2E flow", rgb("A855F7")),
        ("Transaction model", LIME),
        ("ODI", BLUE),
        ("GoldenGate", CYAN),
        ("Qfile", GREEN),
        ("Oracle DB", AMBER),
        ("Lineage", ROSE),
        ("Exception handling", INDIGO),
        ("Activity classification", ORANGE),
    ]
    x, y = 0.75, 1.80
    for idx, (word, fill) in enumerate(words):
        width = max(1.10, min(2.10, 0.40 + 0.11 * len(word)))
        if x + width > 12.55:
            x = 0.75
            y += 0.50
        add_chip(slide, x, y, width, word, fill)
        x += width + 0.15

    add_card(
        slide,
        0.75,
        4.65,
        12.0,
        1.45,
        "Glossary",
        [
            "CDS - the source-to-target control and integration layer",
            "CCAL - the business layer for transaction and activity processing",
            "E2E - end-to-end source to audited business output",
        ],
        accent=BLUE,
        body_size=12,
    )

    add_note(
        slide,
        "Use this slide for search terminology and to anchor the audience on the transaction and activity concepts."
    )
    add_footer(slide)


def create_closing_slide(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, NAVY)
    add_top_band(slide, "Closing summary", page, total)

    add_textbox(slide, 0.65, 0.82, 9.5, 0.45, "What this deck now communicates", font_size=23, bold=True, color=WHITE)
    add_textbox(
        slide,
        0.65,
        1.25,
        8.9,
        0.48,
        "A transaction and activity focused presentation with translated diagrams, clear tables/schemas, and the local style you asked for.",
        font_size=16,
        color=rgb("E2E8F0"),
    )

    roadmap = [
        ("1", "Transaction flow and activity model", CYAN),
        ("2", "Source families, Qfile, and special branches", GREEN),
        ("3", "Tables, schemas, relationships, and controls", AMBER),
        ("4", "Ready for branding or extra screenshots", ROSE),
    ]
    y = 2.0
    for number, label, fill in roadmap:
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.8), Inches(y), Inches(0.48), Inches(0.48))
        dot.fill.solid()
        dot.fill.fore_color.rgb = fill
        dot.line.fill.background()
        tf = dot.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = number
        r.font.name = "Aptos"
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = WHITE

        add_textbox(slide, 1.40, y - 0.03, 9.0, 0.35, label, font_size=16, color=WHITE)
        y += 0.62

    add_card(
        slide,
        8.85,
        1.78,
        3.75,
        2.45,
        "Output path",
        [
            "PP_BOT/data/outputs/CDS_CCAL_Transaction_Activity_Only_Modern_Deck.pptx",
            "This deck is saved locally and can be opened in PowerPoint.",
        ],
        accent=CYAN,
        fill=rgb("0F172A"),
        line=rgb("1E293B"),
        body_size=11,
    )

    add_textbox(
        slide,
        0.65,
        6.0,
        11.2,
        0.30,
        "If needed, I can add brand colors, speaker notes, or replace any translated diagram slide with a different local reference.",
        font_size=12,
        color=rgb("CBD5E1"),
    )

    add_note(
        slide,
        "Finish by summarizing that the deck is now transaction/activity specific and visually closer to the local reference diagram style."
    )
    add_footer(slide)


def build_deck() -> Path:
    prs = Presentation()
    set_wide_layout(prs)

    total = 9
    create_title_slide(prs, 1, total)
    create_scope_slide(prs, 2, total)
    create_source_landscape_slide(prs, 3, total)
    create_reference_style_flow_slide(prs, 4, total)
    create_activity_model_slide(prs, 5, total)
    create_schema_relationship_slide(prs, 6, total)
    create_qfile_slide(prs, 7, total)
    create_orchestration_slide(prs, 8, total)
    create_controls_slide(prs, 9, total)
    # closing slide omitted to keep the deck focused and compact; the prior slide already captures the close.
    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    return OUTPUT_PATH


def main() -> None:
    output = build_deck()
    print(f"Wrote {output}")


if __name__ == "__main__":
    main()
