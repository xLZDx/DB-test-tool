from __future__ import annotations

from pathlib import Path
from typing import Iterable, Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_PATH = Path("PP_BOT/data/outputs/CCAL_Balances_Comprehensive_Modern_Deck.pptx")

# Palette
NAVY = RGBColor(15, 23, 42)
SLATE = RGBColor(51, 65, 85)
MUTED = RGBColor(100, 116, 139)
TEXT = RGBColor(15, 23, 42)
WHITE = RGBColor(255, 255, 255)
SOFT_BG = RGBColor(247, 250, 252)
CARD = RGBColor(255, 255, 255)
BORDER = RGBColor(226, 232, 240)
BLUE = RGBColor(37, 99, 235)
CYAN = RGBColor(6, 182, 212)
GREEN = RGBColor(16, 185, 129)
AMBER = RGBColor(245, 158, 11)
ROSE = RGBColor(236, 72, 153)
INDIGO = RGBColor(99, 102, 241)
LIME = RGBColor(132, 204, 22)
ORANGE = RGBColor(249, 115, 22)
PALE_BLUE = RGBColor(219, 234, 254)
PALE_CYAN = RGBColor(207, 250, 254)
PALE_GREEN = RGBColor(220, 252, 231)
PALE_AMBER = RGBColor(254, 243, 199)
PALE_ROSE = RGBColor(252, 231, 243)
PALE_INDIGO = RGBColor(224, 231, 255)


def rgb(value: str) -> RGBColor:
    value = value.strip().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def set_wide_layout(prs: Presentation) -> None:
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)


def add_full_bg(slide, color: RGBColor) -> None:
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()


def add_top_band(slide, section: str, page: int, total: int) -> None:
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.48))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()

    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0.48), Inches(13.333), Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = CYAN
    accent.line.fill.background()

    label = slide.shapes.add_textbox(Inches(0.45), Inches(0.08), Inches(8.8), Inches(0.22))
    tf = label.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = section
    r.font.name = "Aptos"
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = WHITE

    page_box = slide.shapes.add_textbox(Inches(11.9), Inches(0.08), Inches(1.0), Inches(0.22))
    tf = page_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{page:02d}/{total:02d}"
    r.font.name = "Aptos"
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = WHITE


def add_footer(slide, text: str = "Source baseline: local balance brief, flow images, SQL artifacts, and reference diagrams") -> None:
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(7.14), Inches(13.333), Inches(0.36))
    line.fill.solid()
    line.fill.fore_color.rgb = rgb("EEF2FF")
    line.line.fill.background()

    left = slide.shapes.add_textbox(Inches(0.5), Inches(7.18), Inches(10.4), Inches(0.18))
    tf = left.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(9)
    r.font.color.rgb = MUTED


def add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    font_size: int = 18,
    color: RGBColor = TEXT,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    font_name: str = "Aptos",
) -> any:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font_name
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_multiline(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    lines: Sequence[str],
    font_size: int = 14,
    color: RGBColor = TEXT,
    bullet: bool = False,
    font_name: str = "Aptos",
) -> any:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if bullet:
            p.bullet = True
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line
        r.font.name = font_name
        r.font.size = Pt(font_size)
        r.font.color.rgb = color
    return box


def add_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    body: Iterable[str] | str,
    accent: RGBColor = BLUE,
    fill: RGBColor = CARD,
    line: RGBColor = BORDER,
    title_size: int = 16,
    body_size: int = 12,
) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.2)

    stripe = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.10))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()

    add_textbox(slide, left + 0.18, top + 0.16, width - 0.36, 0.32, title, font_size=title_size, color=TEXT, bold=True)

    if isinstance(body, str):
        body_lines = [body]
    else:
        body_lines = list(body)

    add_multiline(slide, left + 0.18, top + 0.56, width - 0.36, height - 0.70, body_lines, font_size=body_size, color=SLATE)


def add_chip(slide, left: float, top: float, width: float, text: str, fill: RGBColor, color: RGBColor = WHITE) -> None:
    chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.34))
    chip.fill.solid()
    chip.fill.fore_color.rgb = fill
    chip.line.fill.background()
    tf = chip.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = color


def add_arrow(slide, left: float, top: float, width: float, height: float, fill: RGBColor = BLUE) -> None:
    arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(left), Inches(top), Inches(width), Inches(height))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = fill
    arrow.line.fill.background()


def add_line(slide, left: float, top: float, width: float, height: float, color: RGBColor = BORDER) -> None:
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def add_note(slide, text: str) -> None:
    try:
        slide.notes_slide.notes_text_frame.text = text
    except Exception:
        pass


def title_slide(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, NAVY)

    for left, top, width, height, color in [
        (9.00, 0.75, 3.05, 0.18, CYAN),
        (10.05, 1.05, 1.95, 0.18, GREEN),
        (8.55, 1.35, 3.50, 0.18, AMBER),
    ]:
        s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()

    add_textbox(slide, 0.62, 0.42, 7.0, 0.36, "CCAL Balances", font_size=14, color=CYAN, bold=True)
    add_textbox(slide, 0.62, 1.00, 7.1, 1.30, "Comprehensive\nBalance Processing Deck", font_size=34, color=WHITE, bold=True)
    add_textbox(
        slide,
        0.65,
        2.38,
        7.2,
        0.70,
        "A one-hour, speaker-noted review of CCAL balances, EOD and intraday processing,\nsource files, journals, markers, tables, and operating controls.",
        font_size=16,
        color=rgb("E2E8F0"),
    )

    chips = [
        ("Balances", BLUE),
        ("EOD", CYAN),
        ("Intraday", GREEN),
        ("Statement", AMBER),
        ("Cashheld", ROSE),
        ("ACTBALWK", INDIGO),
        ("RJ Bank", ORANGE),
        ("Shadow", rgb("0EA5E9")),
        ("Available", rgb("14B8A6")),
        ("Liquidation", rgb("A855F7")),
        ("AR_ID", LIME),
    ]
    positions = [
        (0.65, 3.55, 1.05),
        (1.78, 3.55, 0.75),
        (2.62, 3.55, 1.00),
        (3.70, 3.55, 1.12),
        (4.92, 3.55, 1.15),
        (6.15, 3.55, 1.38),
        (0.65, 3.98, 1.10),
        (1.86, 3.98, 0.95),
        (2.90, 3.98, 1.10),
        (4.08, 3.98, 1.22),
        (5.38, 3.98, 0.95),
    ]
    for (label, fill), (x, y, w) in zip(chips, positions):
        add_chip(slide, x, y, w, label, fill)

    right = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.45), Inches(0.82), Inches(4.05), Inches(5.75))
    right.fill.solid()
    right.fill.fore_color.rgb = rgb("0B1220")
    right.line.color.rgb = rgb("1E293B")
    add_textbox(slide, 8.80, 1.00, 3.0, 0.34, "What this meeting covers", font_size=16, color=CYAN, bold=True)

    steps = [
        ("Model", "AR_ID, summary vs sub-balance,\nrow layout, snapshot cadence", BLUE),
        ("Families", "Liquidation, statement,\nACTBALWK, Cashheld, RJ Bank", GREEN),
        ("Mechanics", "markers, journals, retries,\nEOD and intraday jobs", AMBER),
        ("Controls", "exceptions, reconciliation,\nstate transitions, clean-up", ROSE),
    ]
    for i, (title, subtitle, fill) in enumerate(steps):
        y = 1.45 + i * 0.90
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.85), Inches(y), Inches(2.95), Inches(0.56))
        box.fill.solid()
        box.fill.fore_color.rgb = fill
        box.line.fill.background()
        tf = box.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = f"{title}\n"
        r.font.name = "Aptos"
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = WHITE
        r2 = p.add_run()
        r2.text = subtitle
        r2.font.name = "Aptos"
        r2.font.size = Pt(9)
        r2.font.color.rgb = WHITE

    for y in [2.05, 2.95, 3.85]:
        add_arrow(slide, 10.08, y, 0.46, 0.22, fill=rgb("334155"))

    add_textbox(
        slide,
        8.78,
        5.20,
        3.10,
        0.55,
        "Designed to be a polished executive/technical review deck,\nnot a raw notes dump.",
        font_size=11,
        color=rgb("CBD5E1"),
    )
    add_note(
        slide,
        "Open with the purpose of the meeting: this deck focuses on balances only, the operating model behind them, and the practical jobs and tables used to move data into CCAL."
    )
    add_footer(slide)


def scope_slide(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, SOFT_BG)
    add_top_band(slide, "Scope and meeting agenda", page, total)
    add_textbox(slide, 0.55, 0.72, 8.8, 0.42, "Balance scope and the operating model", font_size=23, bold=True)
    add_textbox(
        slide,
        0.55,
        1.10,
        11.2,
        0.48,
        "This deck is intentionally balance-specific: no positions deep dive, no transaction deep dive, and no CCAL overview outside the balance context.",
        font_size=15,
        color=MUTED,
    )

    cards = [
        ("1  Balance landscape", "AR_ID, summary and sub-balances,\nEOD vs intraday, detail vs summary", CYAN),
        ("2  Core families", "Liquidation, statement, ACTBALWK,\nCashheld, RJ Bank, Shadow, Available", BLUE),
        ("3  Processing model", "Markers, journals, staging,\nload jobs, clean-up, reconciliation", GREEN),
        ("4  1-hour meeting", "Designed with speaker notes so each slide can carry a 4-6 minute talk track", ROSE),
    ]
    xs = [0.55, 3.75, 6.95, 10.15]
    for (title, body, accent), x in zip(cards, xs):
        add_card(slide, x, 1.80, 2.65, 2.15, title, body.split("\n"), accent=accent, body_size=12, title_size=14)

    add_card(
        slide,
        0.55,
        4.35,
        12.2,
        1.35,
        "Suggested flow for the meeting",
        [
            "Start with the balance model and why AR_ID is the anchor.",
            "Walk through each balance family and its data path.",
            "Use the diagram slide to tie the families together.",
            "Close with BMC/ODI jobs, markers, and the control plane.",
        ],
        accent=INDIGO,
        body_size=13,
    )

    add_note(
        slide,
        "Use this slide to set the agenda and explain that the goal is to cover balance processing families and operational behavior, not to restate the transaction architecture."
    )
    add_footer(slide)


def balance_model_slide(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, WHITE)
    add_top_band(slide, "Balance model at a glance", page, total)
    add_textbox(slide, 0.55, 0.72, 9.2, 0.42, "How CCAL balances are represented", font_size=23, bold=True)

    hub = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(5.00), Inches(2.00), Inches(3.0), Inches(2.1))
    hub.fill.solid()
    hub.fill.fore_color.rgb = NAVY
    hub.line.fill.background()
    tf = hub.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Brokerage\nAR_ID"
    r.font.name = "Aptos"
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = WHITE

    add_card(slide, 0.60, 1.05, 3.15, 1.15, "Row model", ["All balances are stored against Brokerage AR_ID.", "Main and sub-balances live as columns in the same row."], accent=CYAN, body_size=12)
    add_card(slide, 0.60, 2.70, 3.15, 1.15, "What drives balances", ["Balance values can be derived from sources, journals, or position data.", "Security-dependent values are derived in the service layer."], accent=GREEN, body_size=12)
    add_card(slide, 0.60, 4.35, 3.15, 1.15, "Snapshot cadence", ["Most families use an end-of-day snapshot.", "Some families maintain intraday state and EOD freeze points."], accent=AMBER, body_size=12)

    add_card(slide, 9.55, 1.05, 3.15, 1.15, "Detail vs summary", ["Some balances keep separate detail tables (for example deposit holds).", "Summary values are kept in CCAL_BAL."], accent=ROSE, body_size=12)
    add_card(slide, 9.55, 2.70, 3.15, 1.15, "Separate balance tables", ["Some balance families are split into dedicated balance detail tables.", "Sweep and hold information use separate structures."], accent=INDIGO, body_size=12)
    add_card(slide, 9.55, 4.35, 3.15, 1.15, "Operational view", ["The model is designed for auditability, source traceability, and reprocessing.", "Every family has its own job, table, and control path."], accent=BLUE, body_size=12)

    add_arrow(slide, 3.85, 1.55, 1.00, 0.20, fill=rgb("CBD5E1"))
    add_arrow(slide, 3.85, 3.20, 1.00, 0.20, fill=rgb("CBD5E1"))
    add_arrow(slide, 3.85, 4.85, 1.00, 0.20, fill=rgb("CBD5E1"))
    add_arrow(slide, 8.10, 1.55, 1.00, 0.20, fill=rgb("CBD5E1"))
    add_arrow(slide, 8.10, 3.20, 1.00, 0.20, fill=rgb("CBD5E1"))
    add_arrow(slide, 8.10, 4.85, 1.00, 0.20, fill=rgb("CBD5E1"))

    add_textbox(
        slide,
        0.70,
        6.00,
        12.0,
        0.26,
        "The model is consistent: AR_ID is the anchor, balance families are columns/structures, and the cadence is EOD or intraday depending on the source.",
        font_size=12,
        color=MUTED,
    )

    add_note(
        slide,
        "Explain the data model. Emphasize that CCAL balances are not a single monolith; each family uses the same basic anchor, but different source tables, timing, and exception behavior."
    )
    add_footer(slide)


def flow_overview_slide(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, SOFT_BG)
    add_top_band(slide, "Balance flow overview", page, total)
    add_textbox(slide, 0.55, 0.72, 10.0, 0.42, "Three dominant paths shown as a modern translated diagram", font_size=23, bold=True)
    add_textbox(
        slide,
        0.55,
        1.10,
        11.3,
        0.32,
        "Left-to-right flow: source -> journal or file intake -> staging -> target/summary -> consumer snapshot",
        font_size=13,
        color=MUTED,
    )

    panels = [
        ("EOD snapshot path", ["POS", "liquidation", "statement"], CYAN, PALE_CYAN, 0.55),
        ("Intraday journal path", ["ACTBALWK", "Cashheld", "RJ Bank"], GREEN, PALE_GREEN, 4.38),
        ("File / statement path", ["STMTINFO", "Shadow statement", "merge into CCAL_BAL"], AMBER, PALE_AMBER, 8.21),
    ]
    for title, items, accent, fill, x in panels:
        add_card(slide, x, 1.75, 4.15, 4.80, title, items + [" ", ""] if False else items, accent=accent, fill=fill, body_size=12, title_size=15)

    # diagram contents inside each panel
    # panel 1: EOD snapshot
    p1x = 0.75
    add_chip(slide, p1x, 2.25, 1.1, "POS", CYAN)
    add_arrow(slide, p1x + 1.25, 2.36, 0.35, 0.16, fill=rgb("CBD5E1"))
    add_chip(slide, p1x + 1.65, 2.25, 1.45, "EOD stage", BLUE)
    add_arrow(slide, p1x + 3.25, 2.36, 0.35, 0.16, fill=rgb("CBD5E1"))
    add_chip(slide, p1x, 3.05, 1.1, "Liquidation", GREEN)
    add_chip(slide, p1x + 1.35, 3.05, 1.3, "Statement", AMBER)
    add_chip(slide, p1x + 2.85, 3.05, 1.0, "CCAL_BAL", ROSE)
    add_line(slide, p1x + 1.15, 2.58, 0.02, 0.52, color=rgb("CBD5E1"))
    add_line(slide, p1x + 2.60, 2.58, 0.02, 0.52, color=rgb("CBD5E1"))

    # panel 2: intraday
    p2x = 4.58
    add_chip(slide, p2x, 2.25, 1.35, "Journal", CYAN)
    add_arrow(slide, p2x + 1.55, 2.36, 0.35, 0.16, fill=rgb("CBD5E1"))
    add_chip(slide, p2x + 1.95, 2.25, 1.35, "ID_STG", BLUE)
    add_arrow(slide, p2x + 3.45, 2.36, 0.35, 0.16, fill=rgb("CBD5E1"))
    add_chip(slide, p2x + 3.85, 2.25, 0.95, "CCAL_BAL", GREEN)
    add_chip(slide, p2x, 3.05, 1.25, "ACTBALWK", ROSE)
    add_chip(slide, p2x + 1.35, 3.05, 1.1, "Cashheld", INDIGO)
    add_chip(slide, p2x + 2.65, 3.05, 1.15, "RJ Bank", AMBER)
    add_line(slide, p2x + 1.40, 2.58, 0.02, 0.52, color=rgb("CBD5E1"))
    add_line(slide, p2x + 2.80, 2.58, 0.02, 0.52, color=rgb("CBD5E1"))
    add_line(slide, p2x + 4.10, 2.58, 0.02, 0.52, color=rgb("CBD5E1"))

    # panel 3: file/statements
    p3x = 8.41
    add_chip(slide, p3x, 2.25, 1.10, "STMTINFO", CYAN)
    add_arrow(slide, p3x + 1.25, 2.36, 0.35, 0.16, fill=rgb("CBD5E1"))
    add_chip(slide, p3x + 1.65, 2.25, 1.35, "external table", BLUE)
    add_arrow(slide, p3x + 3.15, 2.36, 0.35, 0.16, fill=rgb("CBD5E1"))
    add_chip(slide, p3x, 3.05, 1.15, "EOD stage", GREEN)
    add_chip(slide, p3x + 1.30, 3.05, 1.35, "summary merge", AMBER)
    add_chip(slide, p3x + 2.80, 3.05, 1.10, "CCAL_BAL", ROSE)
    add_line(slide, p3x + 1.25, 2.58, 0.02, 0.52, color=rgb("CBD5E1"))
    add_line(slide, p3x + 2.60, 2.58, 0.02, 0.52, color=rgb("CBD5E1"))

    # bottom callout spanning all panels
    add_card(
        slide,
        0.65,
        5.20,
        12.0,
        1.05,
        "Common pattern",
        [
            "Each family has a source-specific intake method but ultimately lands in a controlled CCAL balance structure.",
            "The timing, snapshot rules, and exception behavior vary; the architecture stays the same.",
        ],
        accent=INDIGO,
        body_size=12,
    )

    add_note(
        slide,
        "Explain this slide as the balance architecture map. Tell the audience that the three lanes are the main patterns used throughout the rest of the deck: EOD snapshot, intraday journaled, and file/statement driven."
    )
    add_footer(slide)


def liquidation_slide(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, WHITE)
    add_top_band(slide, "Liquidation balances", page, total)
    add_textbox(slide, 0.55, 0.72, 8.8, 0.42, "Liquidation balances are derived from position data", font_size=23, bold=True)

    add_card(
        slide,
        0.60,
        1.60,
        4.0,
        1.65,
        "What it is",
        [
            "EOD only; no intraday value.",
            "Derived from CCAL_OWNER.POS.",
            "LQD_VAL is the sum of liquidation sub-balances.",
        ],
        accent=BLUE,
        body_size=12,
    )
    add_card(
        slide,
        4.95,
        1.60,
        3.8,
        1.65,
        "Stage table",
        [
            "CCAL_BAL_STG_OWNER.CCAL_POS_TO_BAL_EOD_STG",
            "Loaded from POS during CCAL batch after source-specific POS load completes.",
        ],
        accent=GREEN,
        body_size=12,
    )
    add_card(
        slide,
        9.05,
        1.60,
        3.7,
        1.65,
        "Target table",
        [
            "CCAL_BAL_OWNER.CCAL_BAL",
            "Column: LQD_VAL",
            "Updated by CCAL_BAL_LOAD_BALANCES_EOD.",
        ],
        accent=ROSE,
        body_size=12,
    )

    sub_cols = [
        "TD_SCR_HELD_CASH",
        "TD_SCR_HELD_MRGN",
        "TD_CASHLESS_STK_OPT",
        "TD_CASH_BAL",
        "TD_MRGN_CASH_BAL",
        "TD_CIP_BAL",
        "TD_BDP_BAL",
        "TD_RJ_BNK_CHK_BAL",
        "TD_RJ_BNK_CD",
    ]
    add_card(
        slide,
        0.60,
        3.60,
        12.15,
        1.05,
        "Liquidation sub balances",
        sub_cols,
        accent=AMBER,
        body_size=11,
    )

    add_card(
        slide,
        0.60,
        4.95,
        5.95,
        1.25,
        "Scheduling notes",
        [
            "CCAL_BAL_LOAD_BALANCES_POS_TO_BAL_EOD_STG loads the EOD staging table.",
            "No dedicated marker is needed; the job runs after POS load for the relevant sources completes.",
        ],
        accent=INDIGO,
        body_size=11,
    )
    add_card(
        slide,
        6.80,
        4.95,
        5.95,
        1.25,
        "Speaker cue",
        [
            "Talk through how one balance value is the aggregate of multiple sub-balances.",
            "The important concept is the EOD snapshot and the source-to-target traceability.",
        ],
        accent=CYAN,
        body_size=11,
    )

    add_note(
        slide,
        "Use this slide to explain the liquidation family. Emphasize that liquidation balances are position-derived, EOD-only, and aggregated from multiple sub-balances into LQD_VAL."
    )
    add_footer(slide)


def statement_slide(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, SOFT_BG)
    add_top_band(slide, "Statement balances", page, total)
    add_textbox(slide, 0.55, 0.72, 8.8, 0.42, "Statement balances are driven by marker-based EOD ingestion", font_size=23, bold=True)

    add_card(
        slide,
        0.60,
        1.55,
        3.6,
        1.7,
        "Current pattern",
        [
            "Consume statement data when the marker arrives.",
            "Load source data into <Source>EOD_STG.",
            "Merge into CCAL_BAL for the EOD date.",
        ],
        accent=CYAN,
        body_size=12,
    )
    add_card(
        slide,
        4.45,
        1.55,
        4.1,
        1.7,
        "Tables and properties",
        [
            "BAL_STMT_ID_DT",
            "BAL_SHDW_STMT_ID_DT",
            "STMTOTFLEOD marker group",
            "BMC / ODI orchestration",
        ],
        accent=GREEN,
        body_size=12,
    )
    add_card(
        slide,
        8.80,
        1.55,
        3.95,
        1.7,
        "Target behavior",
        [
            "Once the data is in CCAL_BAL for the EOD date, the balances are merged for ID date support.",
            "Shadow and brokerage statement processing share the same design pattern.",
        ],
        accent=AMBER,
        body_size=12,
    )

    add_card(
        slide,
        0.60,
        3.65,
        12.15,
        1.65,
        "ODI objects and jobs",
        [
            "CCAL_PKG_LOAD_BALANCES_STMT_EOD",
            "CCAL_PKG_EXECUTE_LOAD_SHDW_STMT_BAL_TGT",
            "CCAL_PKG_LOAD_BALANCES_SHDW_STMT_EOD",
            "CCAL_BAL_LOAD_BALANCES_EOD",
            "CCAL_BAL_PKG_UPD_EOD_TO_ID",
        ],
        accent=ROSE,
        body_size=12,
    )

    add_card(
        slide,
        0.60,
        5.55,
        12.15,
        0.95,
        "Meeting talk track",
        [
            "Explain that statement balances are date-sensitive, marker-driven, and split between current prod and newer file formats.",
        ],
        accent=INDIGO,
        body_size=11,
    )

    add_note(
        slide,
        "Walk through the statement balance lifecycle. Highlight the marker-driven load, the EOD staging table, the merge to CCAL_BAL, and the fact that statement balances support both EOD and ID views."
    )
    add_footer(slide)


def actbalwk_slide(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, WHITE)
    add_top_band(slide, "ACTBALWK balances", page, total)
    add_textbox(slide, 0.55, 0.72, 9.0, 0.42, "ACTBALWK: the intraday and EOD balance engine", font_size=23, bold=True)

    add_card(
        slide,
        0.60,
        1.55,
        4.15,
        2.20,
        "Balances maintained",
        [
            "Deposits on hold",
            "Unsettled trades",
            "Unsettled trades ACH",
            "Pending dividend checks",
            "Margin buy trades",
            "Pending deposit amount",
        ],
        accent=CYAN,
        body_size=12,
    )
    add_card(
        slide,
        4.95,
        1.55,
        3.85,
        2.20,
        "Source and snapshot",
        [
            "CCAL_REPL_OWNER.SIS_ACTBALWK",
            "Journal table J$SIS_ACTBALWK",
            "Marker MRKACTB / group MRKRACTBALWK",
            "EOD snapshot property BAL_ACTBALWK_ID_DT",
        ],
        accent=GREEN,
        body_size=12,
    )
    add_card(
        slide,
        9.05,
        1.55,
        3.70,
        2.20,
        "Data structures",
        [
            "CCAL_BAL_OWNER.DEP_HOLD_BAL_DTL",
            "CCAL_BAL_OWNER.CCAL_BAL",
            "CCAL_BAL_OWNER.CCAL_ACTBALWK_ID_STG",
            "J$CCAL_BAL / FUTR_DT tables",
        ],
        accent=AMBER,
        body_size=12,
    )

    add_card(
        slide,
        0.60,
        4.15,
        12.15,
        1.85,
        "Processing rules",
        [
            "Retry the same journal up to x times when an AR_ID lookup fails.",
            "Move future-dated records to FUTR_DT to keep the EOD snapshot clean.",
            "Use single-commit writes so the client center view stays internally consistent.",
            "Recalculate deposit hold when the business day flips and recalculate deleted source records.",
        ],
        accent=ROSE,
        body_size=12,
    )

    add_card(
        slide,
        0.60,
        6.18,
        12.15,
        0.42,
        "ODI scenarios and jobs: CCAL_PKG_LOAD_BALANCES_RT_ACTBALWK, CCAL_BAL_LOAD_BALANCES_EOD_DEP_DTL, CCAL_BAL_LOAD_BALANCES_EOD_STG_ACTBALWK, CCAL_PKG_LOAD_BALANCES_ACTBALWK_MRGN_INST, CCAL_PKG_ACTBALWK_ACH_ID_UPD_NBD, CCAL_PKG_LOAD_BALANCES_ACTBALWK_FUTR_DT, CCAL_PKG_LOAD_BALANCES_ACTBALWK_FUTR_DT_CLEANUP, CCAL_BAL_LOAD_BALANCES_EOD",
        accent=INDIGO,
        body_size=10,
    )

    add_note(
        slide,
        "Explain ACTBALWK as the key intraday balance engine. Emphasize retries, future-dated record handling, single-commit consistency, and the split between detail and summary outputs."
    )
    add_footer(slide)


def cashheld_slide(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, SOFT_BG)
    add_top_band(slide, "Cashheld balances", page, total)
    add_textbox(slide, 0.55, 0.72, 8.8, 0.42, "Cashheld balances: detail + summary with a strict snapshot model", font_size=23, bold=True)

    add_card(
        slide,
        0.60,
        1.60,
        4.0,
        1.75,
        "Core behavior",
        [
            "GoldenGate replication from SIS.",
            "Trigger tracks journals for ETL consumption.",
            "Lookup failures retry x times then go to exception handling.",
        ],
        accent=BLUE,
        body_size=12,
    )
    add_card(
        slide,
        4.90,
        1.60,
        4.0,
        1.75,
        "Detail + summary",
        [
            "Detail table: CCAL_BAL_OWNER.CASH_HELD_BAL_DTL",
            "Summary table: CCAL_BAL_OWNER.CCAL_BAL",
            "Summary is account level.",
        ],
        accent=GREEN,
        body_size=12,
    )
    add_card(
        slide,
        9.20,
        1.60,
        3.55,
        1.75,
        "Control points",
        [
            "App property BAL_CASHHELD_ID_DT",
            "Marker job CDSCCALPROD.UX.MKCDSBALCASHHELD",
            "Cleanup job removes expired detail rows.",
        ],
        accent=AMBER,
        body_size=12,
    )

    add_card(
        slide,
        0.60,
        3.75,
        12.15,
        1.65,
        "Why it is important",
        [
            "Cashheld is one of the most sensitive balances because users expect both detail and summary to agree at snapshot time.",
            "The process uses a single commit transaction to keep the client-facing view consistent.",
            "Multithreading is enabled and the thread count is driven dynamically by account number.",
        ],
        accent=ROSE,
        body_size=12,
    )

    add_card(
        slide,
        0.60,
        5.65,
        12.15,
        0.80,
        "Common scenarios: marker received or not received, process all or cap at x journals, move future-dated journals into FUTR_DT for replay and recalculation.",
        accent=INDIGO,
        body_size=11,
    )

    add_note(
        slide,
        "Describe cashheld as a detail and summary balance with strict snapshot timing. Highlight the journal retry rules, single commit, marker behavior, and the cleanup of expired detail rows."
    )
    add_footer(slide)


def rj_bank_slide(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, WHITE)
    add_top_band(slide, "RJ Bank balances", page, total)
    add_textbox(slide, 0.55, 0.72, 9.0, 0.42, "RJ Bank balances: checking and CDs", font_size=23, bold=True)

    add_card(
        slide,
        0.60,
        1.60,
        4.1,
        1.90,
        "Source and replication",
        [
            "SIS_BANKBLNC and SIS_BANKMSTR",
            "Journal table J$SIS_BANKBLNC",
            "Triggers T$SIS_BANKBLNC and T$SIS_BANKMSTR",
        ],
        accent=CYAN,
        body_size=12,
    )
    add_card(
        slide,
        4.95,
        1.60,
        4.05,
        1.90,
        "Intraday flow",
        [
            "Common intraday position table CCAL_STG_OWNER.CCAL_ID_POS_STG",
            "Target DLY_POS_EXTRACT",
            "Merge proc ID_TXN_POS_WFLOW_TGT",
        ],
        accent=GREEN,
        body_size=12,
    )
    add_card(
        slide,
        9.25,
        1.60,
        3.5,
        1.90,
        "Control properties",
        [
            "BAL_RJBNK_BLNC_ID_DT",
            "BAL_RJBNK_CHK_ID_DT",
            "BAL_RJBNK_BLNC_RTR_LIMIT / BAL_RJBNK_CHK_RTR_LIMIT",
        ],
        accent=AMBER,
        body_size=12,
    )

    add_card(
        slide,
        0.60,
        3.85,
        12.15,
        1.65,
        "What the balance family does",
        [
            "Bank checking and bank CD balances are processed as a dedicated intraday balance family.",
            "The flows retry on AR_ID lookup failures, then move unresolved items to source-specific exception tables.",
            "Once the marker arrives, the business date property is flipped and the balance flow closes out the snapshot.",
        ],
        accent=ROSE,
        body_size=12,
    )

    add_card(
        slide,
        0.60,
        5.80,
        12.15,
        0.68,
        "ODI / proc names: CCAL_PKG_LOAD_BALANCES_RT_RJBNKBLNC, CCAL_PKG_LOAD_ID_POS_TXN_WFLOW_TGT, CCAL_PKG_LOAD_ID_REINITIALIE, CCAL_ID_UTIL_PKG.ID_POS_RJBNK_BLNC_STG, CCAL_ID_UTIL_PKG.ID_REINITIALIE",
        accent=INDIGO,
        body_size=10,
    )

    add_note(
        slide,
        "Present RJ Bank as a dedicated intraday balance branch with checking and CD variants. Mention the common position staging table, the target extract, and the retry/exceptions behavior."
    )
    add_footer(slide)


def shadow_statement_slide(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, SOFT_BG)
    add_top_band(slide, "Shadow statement balances", page, total)
    add_textbox(slide, 0.55, 0.72, 10.0, 0.42, "Shadow statement balances and file-transfer path", font_size=23, bold=True)

    add_card(
        slide,
        0.60,
        1.55,
        3.7,
        1.85,
        "Current prod pattern",
        [
            "File: STMTINFO",
            "External table: CDS_STG_OWNER.SIS_STMTINFO_EXT",
            "Fill table: CCAL_FILL_OWNER.SHDW_SIS_STMTINFO",
        ],
        accent=BLUE,
        body_size=12,
    )
    add_card(
        slide,
        4.60,
        1.55,
        4.05,
        1.85,
        "New file pattern",
        [
            "File: STMTINFO with INFO_RECORD_TYPE",
            "External table: CDS_STG_OWNER.SIS_STMTINFO_BAL_EXT",
            "Stage table: CDS_STG_OWNER.SHDW_SIS_STMTINFO_BAL_STG",
        ],
        accent=GREEN,
        body_size=12,
    )
    add_card(
        slide,
        8.95,
        1.55,
        3.35,
        1.85,
        "The common idea",
        [
            "File moves through Nighthawk -> NFS -> ODI load -> CCAL staging.",
            "The BMC job names remain the same even when scenario names change.",
        ],
        accent=AMBER,
        body_size=12,
    )

    flow_y = 4.05
    add_chip(slide, 0.85, flow_y, 1.45, "Nighthawk", CYAN)
    add_arrow(slide, 2.40, flow_y + 0.09, 0.36, 0.16, fill=rgb("CBD5E1"))
    add_chip(slide, 2.85, flow_y, 1.20, "NFS", BLUE)
    add_arrow(slide, 4.15, flow_y + 0.09, 0.36, 0.16, fill=rgb("CBD5E1"))
    add_chip(slide, 4.60, flow_y, 1.70, "External table", GREEN)
    add_arrow(slide, 6.45, flow_y + 0.09, 0.36, 0.16, fill=rgb("CBD5E1"))
    add_chip(slide, 6.90, flow_y, 1.45, "EOD stage", AMBER)
    add_arrow(slide, 8.45, flow_y + 0.09, 0.36, 0.16, fill=rgb("CBD5E1"))
    add_chip(slide, 8.90, flow_y, 1.45, "CCAL_BAL", ROSE)

    add_card(
        slide,
        0.60,
        4.85,
        12.15,
        1.35,
        "Jobs and properties",
        [
            "CCAL_PKG_EXECUTE_LOAD_SHDW_STMT_BAL_TGT",
            "CCAL_PKG_LOAD_BALANCES_SHDW_STMT_EOD",
            "CCAL_BAL_LOAD_BALANCES_EOD",
            "BAL_SHDWSTMT_ID_DT",
        ],
        accent=INDIGO,
        body_size=12,
    )

    add_note(
        slide,
        "Walk through the month-end shadow statement path first, then mention the newer file format and stage table. Make clear that the file-transfer chain is part of the operational design."
    )
    add_footer(slide)


def available_balances_slide(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, WHITE)
    add_top_band(slide, "Available balances", page, total)
    add_textbox(slide, 0.55, 0.72, 9.2, 0.42, "Available balance families and merge pattern", font_size=23, bold=True)

    add_card(
        slide,
        0.60,
        1.55,
        4.1,
        1.95,
        "ETL design",
        [
            "Consume data from ID_STG and EOD_STG.",
            "Merge balances into CCAL_BAL for both ID and EOD.",
            "ID mode can use a GTT before the final merge.",
        ],
        accent=CYAN,
        body_size=12,
    )
    add_card(
        slide,
        4.95,
        1.55,
        4.1,
        1.95,
        "Balance families",
        [
            "Available for withdrawal",
            "Available for trading",
            "Unapplied cash including unsettled trades",
            "RJ Bank deposit program cash",
            "Chargeable cash and margin interest",
        ],
        accent=GREEN,
        body_size=12,
    )
    add_card(
        slide,
        9.30,
        1.55,
        3.35,
        1.95,
        "Why it matters",
        [
            "These are user-facing balances that often drive account availability.",
            "They are merged from multiple source states and can be both ID and EOD.",
        ],
        accent=AMBER,
        body_size=12,
    )

    add_card(
        slide,
        0.60,
        3.85,
        12.15,
        1.85,
        "Balance families listed in the brief",
        [
            "Withdrawal cash / margin",
            "ACH withdrawal cash / margin",
            "Institutional variants",
            "RJ Bank deposit program cash",
            "Chargeable cash debit interest accrued",
            "Chargeable margin debit interest accrued",
            "Unapplied cash with unsettled trades",
        ],
        accent=ROSE,
        body_size=12,
    )

    add_card(
        slide,
        0.60,
        5.95,
        12.15,
        0.52,
        "ODI objects: CCAL_BAL_LOAD_BALANCES_EOD and CCAL_BAL_LOAD_BALANCES_ID_RT",
        accent=INDIGO,
        body_size=11,
    )

    add_note(
        slide,
        "Explain available balances as the consumer-facing family: the business cares about what clients can withdraw or trade, and the pipeline merges ID and EOD sources into CCAL_BAL."
    )
    add_footer(slide)


def control_plane_slide(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, SOFT_BG)
    add_top_band(slide, "Controls, jobs, markers, and exception handling", page, total)
    add_textbox(slide, 0.55, 0.72, 9.4, 0.42, "How CCAL balance processing stays operationally safe", font_size=23, bold=True)

    flow_labels = [
        ("Marker", "journal ready", CYAN),
        ("BMC job", "schedule trigger", BLUE),
        ("ODI", "load / merge", GREEN),
        ("Exception", "retry / reclassify", AMBER),
        ("Property", "advance business day", ROSE),
        ("Monitor", "reconcile", INDIGO),
    ]
    xs = [0.75, 2.80, 4.85, 6.90, 8.95, 11.00]
    for (title, sub, fill), x in zip(flow_labels, xs):
        add_chip(slide, x, 1.70, 1.6, title, fill)
        add_textbox(slide, x - 0.18, 2.05, 1.95, 0.28, sub, font_size=9, color=MUTED, align=PP_ALIGN.CENTER)

    for ax in [2.38, 4.43, 6.48, 8.53, 10.58]:
        add_arrow(slide, ax, 1.84, 0.35, 0.16, fill=rgb("CBD5E1"))

    add_card(
        slide,
        0.60,
        2.65,
        5.9,
        1.85,
        "Common controls",
        [
            "Marker jobs confirm journal readiness.",
            "Application properties keep track of business day / snapshot state.",
            "Source-specific exceptions capture failed AR_ID or lookup issues.",
        ],
        accent=CYAN,
        body_size=12,
    )
    add_card(
        slide,
        6.75,
        2.65,
        5.95,
        1.85,
        "Common artifacts",
        [
            "Markers: MRKACTB, MRKCASHH, STMTOTFLEOD, MRKBANKBLNC",
            "BMC / ODI orchestration jobs",
            "CCAL_UTIL exception tables and retry paths",
        ],
        accent=GREEN,
        body_size=12,
    )

    add_card(
        slide,
        0.60,
        4.75,
        12.15,
        1.45,
        "Job families and properties called out in the brief",
        [
            "CCAL_PKG_LOAD_BALANCES_EOD, CCAL_BAL_LOAD_BALANCES_ID_RT, CCAL_BAL_PKG_UPD_EOD_TO_ID, CCAL_PKG_GEN_ID_UPD_NBD",
            "BAL_ACTBALWK_ID_DT, BAL_CASHHELD_ID_DT, BAL_RJBNK_BLNC_ID_DT, BAL_STMT_ID_DT, BAL_SHDWSTMT_ID_DT",
        ],
        accent=ROSE,
        body_size=11,
    )

    add_card(
        slide,
        0.60,
        6.35,
        12.15,
        0.28,
        "Talk track: all families use the same operating pattern - marker arrives -> jobs run -> balance state advances -> exceptions are logged -> consumers see a clean snapshot.",
        accent=INDIGO,
        body_size=10,
    )

    add_note(
        slide,
        "This is the control-plane slide. Focus on markers, BMC jobs, ODI packages, application properties, retries, and exception handling as the common operational pattern behind all balance families."
    )
    add_footer(slide)


def closing_slide(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, NAVY)
    add_top_band(slide, "Closing summary and next steps", page, total)

    add_textbox(slide, 0.65, 0.82, 9.5, 0.45, "What the audience should remember", font_size=23, bold=True, color=WHITE)
    add_textbox(
        slide,
        0.65,
        1.25,
        9.2,
        0.48,
        "CCAL balances are a set of related families, each with its own cadence, tables, markers, and operating rules - but all of them use the same trust model: traceable, auditable, and snapshot-driven.",
        font_size=16,
        color=rgb("E2E8F0"),
    )

    roadmap = [
        ("1", "Balance model and AR_ID anchor", CYAN),
        ("2", "Family-by-family differences", GREEN),
        ("3", "Jobs, markers, and properties", AMBER),
        ("4", "Open questions for the 1-hour review", ROSE),
    ]
    y = 2.0
    for number, label, fill in roadmap:
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.8), Inches(y), Inches(0.48), Inches(0.48))
        dot.fill.solid()
        dot.fill.fore_color.rgb = fill
        dot.line.fill.background()
        tf = dot.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = number
        r.font.name = "Aptos"
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = WHITE
        add_textbox(slide, 1.40, y - 0.03, 9.0, 0.35, label, font_size=16, color=WHITE)
        y += 0.62

    add_card(
        slide,
        8.85,
        1.78,
        3.75,
        2.45,
        "Output path",
        [
            "PP_BOT/data/outputs/CCAL_Balances_Comprehensive_Modern_Deck.pptx",
            "The deck includes speaker notes across all slides.",
        ],
        accent=CYAN,
        fill=rgb("0F172A"),
        line=rgb("1E293B"),
        body_size=11,
    )

    add_textbox(
        slide,
        0.65,
        6.0,
        11.5,
        0.30,
        "If you want, I can add source screenshots, more branding, or a PDF export step after the PPTX is finalized.",
        font_size=12,
        color=rgb("CBD5E1"),
    )

    add_note(
        slide,
        "Close with the core message: balances are a family of controlled, auditable, snapshot-driven processes. Invite discussion on any family that needs deeper detail or a dedicated appendix."
    )
    add_footer(slide)


def build_deck() -> Path:
    prs = Presentation()
    set_wide_layout(prs)

    total = 13
    title_slide(prs, 1, total)
    scope_slide(prs, 2, total)
    balance_model_slide(prs, 3, total)
    flow_overview_slide(prs, 4, total)
    liquidation_slide(prs, 5, total)
    statement_slide(prs, 6, total)
    actbalwk_slide(prs, 7, total)
    cashheld_slide(prs, 8, total)
    rj_bank_slide(prs, 9, total)
    shadow_statement_slide(prs, 10, total)
    available_balances_slide(prs, 11, total)
    control_plane_slide(prs, 12, total)
    closing_slide(prs, 13, total)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    return OUTPUT_PATH


def main() -> None:
    output = build_deck()
    print(f"Wrote {output}")


if __name__ == "__main__":
    main()
